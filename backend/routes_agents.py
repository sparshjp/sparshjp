"""Kairos AI Engine v4 — Full Access, Subagents, Expanded Context.
Parallel tool execution, compound tools (scaffold_module, create_page),
auto-restart, compressed results, web search, screenshot, subagents, image gen."""
from fastapi import APIRouter, HTTPException, UploadFile, File
from datetime import datetime, timezone
import uuid
import os
import json
import glob
import subprocess
import shlex
import asyncio
import httpx
import re
import logging
import base64
from kairos_subagents import call_subagent, generate_image as gen_image

router = APIRouter(prefix="/agents", tags=["AI Engine"])

EMERGENT_KEY = None
GROQ_KEY = ""
OPENROUTER_KEY = ""
ANTHROPIC_API_KEY = ""
OPENAI_API_KEY = ""
CEREBRAS_KEY = ""
HUGGINGFACE_KEY = ""
db = None

def set_config(key, database):
    global EMERGENT_KEY, db, GROQ_KEY, OPENROUTER_KEY, ANTHROPIC_API_KEY, OPENAI_API_KEY, CEREBRAS_KEY, HUGGINGFACE_KEY
    EMERGENT_KEY = key
    db = database
    GROQ_KEY = os.environ.get("GROQ_API_KEY", "")
    OPENROUTER_KEY = os.environ.get("OPENROUTER_API_KEY", "")
    ANTHROPIC_API_KEY = os.environ.get("ANTHROPIC_API_KEY", "")
    OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY", "")
    CEREBRAS_KEY = os.environ.get("CEREBRAS_API_KEY", "")
    HUGGINGFACE_KEY = os.environ.get("HUGGINGFACE_API_KEY", "")

# ══════════════════════════════════════════════════════════
# MULTI-PROVIDER LLM CLIENT (Claude → Gemini → GPT-5 → Groq → OpenRouter)
# ══════════════════════════════════════════════════════════

def _call_groq_sync(system: str, messages: list) -> str:
    from groq import Groq
    client = Groq(api_key=GROQ_KEY)
    msgs = [{"role": "system", "content": system}]
    msgs.extend(messages)
    response = client.chat.completions.create(
        model="llama-3.3-70b-versatile", messages=msgs, max_tokens=16000, temperature=0.3,
    )
    return response.choices[0].message.content

def _call_openrouter_sync(system: str, messages: list) -> str:
    from openai import OpenAI
    client = OpenAI(
        base_url="https://openrouter.ai/api/v1", api_key=OPENROUTER_KEY,
        default_headers={"HTTP-Referer": "https://kairos-erp.app", "X-Title": "Kairos AI Engine"},
    )
    msgs = [{"role": "system", "content": system}]
    msgs.extend(messages)
    response = client.chat.completions.create(
        model="openrouter/auto", messages=msgs, max_tokens=16000, temperature=0.3,
    )
    return response.choices[0].message.content

async def _call_claude(system: str, messages: list) -> str:
    from emergentintegrations.llm.chat import LlmChat, UserMessage
    chat = LlmChat(
        api_key=EMERGENT_KEY, session_id=f"engine-{uuid.uuid4()}", system_message=system,
    ).with_model("anthropic", "claude-sonnet-4-5-20250929")
    combined = "\n".join([f"[{m['role'].upper()}]: {m['content']}" for m in messages])
    return await chat.send_message(UserMessage(text=combined))

async def _call_gemini(system: str, messages: list) -> str:
    from emergentintegrations.llm.chat import LlmChat, UserMessage
    chat = LlmChat(
        api_key=EMERGENT_KEY, session_id=f"engine-gemini-{uuid.uuid4()}", system_message=system,
    ).with_model("gemini", "gemini-3-flash-preview")
    combined = "\n".join([f"[{m['role'].upper()}]: {m['content']}" for m in messages])
    combined += "\n\nIMPORTANT FORMAT: Tool calls must use ```TOOL_CALL code blocks. When done, write DONE: summary."
    resp = await chat.send_message(UserMessage(text=combined))
    if resp is None or (isinstance(resp, str) and len(resp.strip()) == 0):
        raise Exception("Gemini returned empty response — falling back to next provider")
    return resp

async def _call_gpt5(system: str, messages: list) -> str:
    from emergentintegrations.llm.chat import LlmChat, UserMessage
    chat = LlmChat(
        api_key=EMERGENT_KEY, session_id=f"engine-gpt5-{uuid.uuid4()}", system_message=system,
    ).with_model("openai", "gpt-5")
    combined = "\n".join([f"[{m['role'].upper()}]: {m['content']}" for m in messages])
    return await chat.send_message(UserMessage(text=combined))


async def _call_claude_direct(system: str, messages: list) -> str:
    """Call Claude directly using user's own Anthropic API key (zero Emergent credits)."""
    import anthropic
    client = anthropic.AsyncAnthropic(api_key=ANTHROPIC_API_KEY)
    msgs = []
    for m in messages:
        msgs.append({"role": m["role"], "content": m["content"]})
    response = await client.messages.create(
        model="claude-sonnet-4-5-20250929",
        max_tokens=16000,
        system=system,
        messages=msgs,
        temperature=0.3,
    )
    return response.content[0].text


async def _call_gpt_direct(system: str, messages: list) -> str:
    """Call OpenAI directly using user's own API key (zero Emergent credits)."""
    from openai import AsyncOpenAI
    client = AsyncOpenAI(api_key=OPENAI_API_KEY)
    msgs = [{"role": "system", "content": system}]
    msgs.extend(messages)
    response = await client.chat.completions.create(
        model="gpt-4o",
        messages=msgs,
        max_tokens=16000,
        temperature=0.3,
    )
    return response.choices[0].message.content


def _call_cerebras_sync(system: str, messages: list) -> str:
    """Call Cerebras free tier — Llama 3.3 70B at ~2000 tok/s. Free: 1M tokens/day."""
    from cerebras.cloud.sdk import Cerebras
    client = Cerebras(api_key=CEREBRAS_KEY)
    msgs = [{"role": "system", "content": system}]
    msgs.extend(messages)
    response = client.chat.completions.create(
        model="llama-3.3-70b", messages=msgs, max_tokens=16000, temperature=0.3,
    )
    return response.choices[0].message.content

def _call_huggingface_sync(system: str, messages: list) -> str:
    """Call HuggingFace free inference — Qwen 2.5 Coder 32B. Free tier available."""
    from huggingface_hub import InferenceClient
    client = InferenceClient(api_key=HUGGINGFACE_KEY)
    msgs = [{"role": "system", "content": system}]
    msgs.extend(messages)
    response = client.chat_completion(
        model="Qwen/Qwen2.5-Coder-32B-Instruct", messages=msgs, max_tokens=16000, temperature=0.3,
    )
    return response.choices[0].message.content

# Track recent provider failures for smart routing
_provider_failures = {}  # {"groq": [timestamp, ...], ...}
_FAILURE_WINDOW = 300  # 5 minutes — skip provider if failed recently

def _should_skip_provider(provider: str) -> bool:
    """Skip a provider if it has failed 2+ times in the last 5 minutes."""
    now = datetime.now(timezone.utc).timestamp()
    failures = _provider_failures.get(provider, [])
    # Clean old failures
    recent = [t for t in failures if now - t < _FAILURE_WINDOW]
    _provider_failures[provider] = recent
    return len(recent) >= 2

def _record_failure(provider: str):
    now = datetime.now(timezone.utc).timestamp()
    _provider_failures.setdefault(provider, []).append(now)

def _clear_failures(provider: str):
    _provider_failures[provider] = []


async def call_llm(system: str, messages: list, preferred: str = "auto") -> tuple:
    PROVIDER_ORDERS = {
        "claude": ["claude_direct", "claude", "gpt_direct", "gemini", "gpt5", "groq", "cerebras", "huggingface", "openrouter"],
        "gemini": ["gemini", "claude_direct", "claude", "gpt_direct", "gpt5", "groq", "cerebras", "huggingface", "openrouter"],
        "gpt5": ["gpt_direct", "gpt5", "claude_direct", "claude", "gemini", "groq", "cerebras", "huggingface", "openrouter"],
        "groq": ["groq", "cerebras", "huggingface", "claude_direct", "claude", "gpt_direct", "gemini", "gpt5", "openrouter"],
        "cerebras": ["cerebras", "groq", "huggingface", "claude_direct", "claude", "gpt_direct", "gemini", "gpt5", "openrouter"],
        "huggingface": ["huggingface", "cerebras", "groq", "claude_direct", "claude", "gpt_direct", "gemini", "gpt5", "openrouter"],
        "openrouter": ["openrouter", "groq", "cerebras", "huggingface", "claude_direct", "claude", "gpt_direct", "gemini", "gpt5"],
    }
    # Smart default: free providers first (zero credits), then direct keys, then Emergent
    default_order = []
    # Free providers first
    if GROQ_KEY:
        default_order.append("groq")
    if CEREBRAS_KEY:
        default_order.append("cerebras")
    if HUGGINGFACE_KEY:
        default_order.append("huggingface")
    # Then user's own paid keys
    if ANTHROPIC_API_KEY:
        default_order.append("claude_direct")
    if OPENAI_API_KEY:
        default_order.append("gpt_direct")
    if OPENROUTER_KEY:
        default_order.append("openrouter")
    # Then Emergent credits as fallback
    default_order.extend(["claude", "gemini", "gpt5"])
    # Deduplicate while preserving order
    seen = set()
    default_order = [x for x in default_order if not (x in seen or seen.add(x))]

    order = PROVIDER_ORDERS.get(preferred, default_order)
    if preferred not in PROVIDER_ORDERS:
        order = default_order
    errors = []
    loop = asyncio.get_event_loop()
    for provider in order:
        if _should_skip_provider(provider):
            errors.append(f"{provider}: skipped (recent failures)")
            continue
        try:
            if provider == "claude_direct" and ANTHROPIC_API_KEY:
                text = await _call_claude_direct(system, messages)
                _clear_failures(provider)
                return text, "claude (direct key)"
            elif provider == "gpt_direct" and OPENAI_API_KEY:
                text = await _call_gpt_direct(system, messages)
                _clear_failures(provider)
                return text, "gpt-4o (direct key)"
            elif provider == "groq" and GROQ_KEY:
                text = await loop.run_in_executor(None, _call_groq_sync, system, messages)
                _clear_failures(provider)
                return text, "groq (free)"
            elif provider == "cerebras" and CEREBRAS_KEY:
                text = await loop.run_in_executor(None, _call_cerebras_sync, system, messages)
                _clear_failures(provider)
                return text, "cerebras (free)"
            elif provider == "huggingface" and HUGGINGFACE_KEY:
                text = await loop.run_in_executor(None, _call_huggingface_sync, system, messages)
                _clear_failures(provider)
                return text, "huggingface (free)"
            elif provider == "openrouter" and OPENROUTER_KEY:
                text = await loop.run_in_executor(None, _call_openrouter_sync, system, messages)
                _clear_failures(provider)
                return text, "openrouter"
            elif provider == "claude" and EMERGENT_KEY:
                text = await _call_claude(system, messages)
                _clear_failures(provider)
                return text, "claude"
            elif provider == "gemini" and EMERGENT_KEY:
                text = await _call_gemini(system, messages)
                _clear_failures(provider)
                return text, "gemini"
            elif provider == "gpt5" and EMERGENT_KEY:
                text = await _call_gpt5(system, messages)
                _clear_failures(provider)
                return text, "gpt5"
        except Exception as e:
            err_msg = str(e)[:200]
            logging.warning(f"AI Engine: {provider} failed: {err_msg}")
            _record_failure(provider)
            errors.append(f"{provider}: {err_msg}")
            continue
    raise Exception(f"All LLM providers failed: {'; '.join(errors)}")

# ══════════════════════════════════════════════════════════
# PATH SAFETY
# ══════════════════════════════════════════════════════════
ALLOWED_DIRS = ["/app/backend", "/app/frontend/src", "/app/frontend/public"]
BLOCKED_PATTERNS = ["node_modules", "__pycache__", ".git", ".emergent"]

def is_safe_path(path):
    for blocked in BLOCKED_PATTERNS:
        if blocked in path:
            return False
    for allowed in ALLOWED_DIRS:
        if path.startswith(allowed):
            return True
    return False

# ══════════════════════════════════════════════════════════
# SYSTEM PROMPT (v3 — speed optimized)
# ══════════════════════════════════════════════════════════
MAX_ITERATIONS = 10

ENGINE_SYSTEM_PROMPT = """You are the Kairos AI Engine v4 — an AUTONOMOUS, senior-level full-stack developer for ABC Ltd IT ERP.
You execute tasks immediately without planning pauses. You think like a principal engineer: plan internally, execute decisively, verify rigorously.

COMPANY: ABC Technologies Pvt. Ltd. | GSTIN: 27AABCA1234P1Z5 | Maharashtra | IT Services
Revenue: INR/USD(84.50)/GBP(106.80) | 8 Projects, 20 Employees, 7 Clients, 10 Vendors
Banks: HDFC(6840000), Axis(2250000), EEFC USD(3042000) | TB: 28142000 (balanced) | 26 CoA

TECH: FastAPI+Motor(MongoDB) backend:8001 | React+Tailwind+Shadcn frontend:3000
Design: Dark #0D1B2A bg, #152236 cards, #1B2D42 borders, #E8EDF2 text, #00d4aa accent

## ERP MODULE MAP (22 Modules — All Interlinked)

### Core Modules & Endpoints
- **Projects** `/api/projects` — GET list, GET /{id}, GET /{id}/timesheets, GET /{id}/transactions, GET /health/dashboard, PUT /{id}/status. DB: `projects`. **MISSING: POST create — ADD IT.**
- **Timesheets** `/api/timesheets` — GET list(?employee_id,week,project_id), POST create, PUT /{id}/approve, PUT /{id}/reject, GET /utilization, GET /consolidation, GET /employees. DB: `timesheets`
- **Approvals** `/api/approvals` — GET /workflows, POST /workflows, GET /requests, POST /requests, POST /requests/{id}/approve, POST /requests/{id}/reject, GET /stats. DB: `approval_workflows`, `approval_requests`
- **Budgets** `/api/budgets` — GET list, POST create, GET /variance, GET /alerts. DB: `budgets`
- **Contracts** `/api/contracts` — GET list, POST create, GET /{id}, GET /alerts/renewals, GET /stats/summary, POST /{id}/milestones/{ms_id}/complete. DB: `contracts`
- **Resources** `/api/resources` — GET /allocations, POST /allocations, GET /bench, GET /utilization, GET /forecast. DB: `resource_allocations`
- **Forex** `/api/forex` — GET /rates, POST /rates/fetch-live, GET /transactions, POST /transactions, POST /transactions/{id}/settle, GET /revaluation. DB: `forex_rates`, `forex_transactions`
- **Billing** `/api/billing` — GET /stats, GET /unbilled, POST /generate-invoice, GET /milestone-invoices, POST /milestone-invoice. DB: `billing_invoices`
- **Documents** `/api/documents` — GET list, POST /upload, GET /download/{id}, DELETE /{id}, GET /categories, GET /stats. DB: `erp_documents`
- **Notifications** `/api/notifications` — GET list, POST create, PUT /{id}/read, PUT /read-all, DELETE /{id}, GET /unread-count, POST /generate-reminders. DB: `notifications`
- **Compliance** `/api/compliance` — GET /frameworks, GET /dashboard, PUT /controls/{fw}/{id}, GET /access-logs, POST /access-logs. DB: `compliance_controls`, `compliance_access_logs`
- **Portal** `/api/portal` — GET /clients, POST /clients, DELETE /clients/{id}, GET /my/dashboard, GET /my/projects, GET /my/invoices. DB: `portal_clients`

### INTER-MODULE LINKING RULES (Event-Driven Pipeline)
When you modify any module, ALWAYS implement these cross-module triggers:

1. **Contract Created → Auto-create Project**: When POST /contracts with milestones, also insert into `projects` collection with matching project_id, client, value, type, milestones.
2. **Contract Milestone Completed → Billing + Notification**: When POST /contracts/{id}/milestones/{ms_id}/complete, insert a record in `billing_invoices` (draft) AND insert a notification ("Milestone X completed, invoice ready").
3. **Timesheet Approved → Billing Queue + Resource Update**: When PUT /timesheets/{id}/approve, mark those entries as billing-ready in the timesheet doc (invoiceable=true). Update resource utilization.
4. **Budget Threshold Exceeded → Approval Request + Notification**: When POST /budgets or budget variance check shows >80% usage, auto-create an approval_request of type "budget_override" AND a notification.
5. **Approval Approved/Rejected → Notification**: When POST /approvals/requests/{id}/approve or reject, insert a notification for the requester.
6. **Resource Allocated → Link to Project Team**: When POST /resources/allocations, update the project's team_names array.
7. **Contract Expiring (<30d) → Notification**: The /contracts/alerts/renewals endpoint feeds into /notifications/generate-reminders.
8. **Invoice Generated → Forex Transaction**: When billing generates an invoice for a non-INR contract, auto-create a forex_transaction.
9. **Document Uploaded → Compliance Log**: When POST /documents/upload, log to compliance_access_logs.

### Frontend Pages (all in /app/frontend/src/pages/)
Projects: ProjectsModule.js — NEEDS "New Project" form with fields: name, client, type(Fixed-Price/T&M/Retainer/etc), pm, value, currency, team_names[], milestones[], billing, duration, status, health
Timesheets: TimesheetsPage.js — NEEDS "New Timesheet" form with fields: employee_id, employee_name, week, week_start, week_end, entries[{project_id, hours, billable, note, rate, currency}], leave_hours, leave_type
Both pages are currently READ-ONLY with no create/edit UI. This is a critical gap.

### DB Schema Quick Reference
- projects: {id, name, client, type, pm, status, health, pct_complete, value_inr, value_usd, currency, billing, duration, team_names[], milestones[{id,name,value,currency,status,date}]}
- timesheets: {id, employee_id, employee_name, week, week_start, week_end, total_hours, status, entries[{project_id,hours,billable,note,rate,currency,ot_hours}], leave_hours, leave_type}
- approval_workflows: {id, name, type, threshold_amount, steps[{role,label}], is_active}
- approval_requests: {id, type, reference_name, amount, requester_name, status, steps[], comments}
- budgets: {id, name, type, department, fiscal_year, line_items[{category,amount,actual}], status}
- contracts: {id, contract_number, title, type, client_name, start_date, end_date, value, currency, billing_type, auto_renew, status, milestones[{id,name,amount,status,invoiced}]}
- resource_allocations: {id, employee_name, project_name, role, allocation_pct, start_date, end_date, billable, bill_rate}
- forex_rates: {base_currency, rates:{USD:x,...}, date, source}
- forex_transactions: {id, type, reference_name, currency, foreign_amount, booking_rate, booking_inr, settlement_rate, settled, forex_gain_loss}
- billing_invoices: {id, project_id, project_name, client, period, entries[], total_amount, status(draft/sent/paid), source(timesheet/milestone)}
- erp_documents: {id, filename, content_type, size, entity_type, entity_id, entity_name, category, uploaded_at, file_path}
- notifications: {id, title, message, type, priority, read, target_roles[], created_at}
- compliance_controls: stored in frameworks dict, compliance_access_logs: {id, user_name, action, resource, ip_address, timestamp}
- portal_clients: {id, client_name, contact_name, email, portal_token, is_active, projects[]}

## REASONING METHODOLOGY (CRITICAL)
Before executing, mentally:
1. **Decompose** — Break task into atomic steps. Identify dependencies.
2. **Risk assess** — What could break? What are the edge cases?
3. **Plan tool calls** — Batch all independent operations into ONE response.
4. **Execute** — Issue tool calls. No planning-only responses.
5. **Verify** — After changes, ALWAYS verify with test_api or verify_deployment.
6. **Self-heal** — If something fails, read logs, diagnose, fix, retry. Never stop at an error.

## DEBUGGING DISCIPLINE
When fixing bugs:
1. **Reproduce first** — Read the file, understand the current state before changing anything.
2. **Trace the chain** — Follow the error from symptom → cause. Don't patch symptoms.
3. **Fix root cause** — One precise change, not shotgun fixes.
4. **Verify fix** — test_api or verify_deployment after every fix.
5. **Regression check** — Ensure the fix doesn't break adjacent functionality.

## TOKEN EFFICIENCY
- Keep tool call args minimal. Don't repeat file content in responses.
- Use patch_file (search/replace) over write_file for targeted edits.
- Use scaffold_module for new backend modules (1 tool = 5+ manual steps).
- Compress your responses. Be terse. Code speaks louder than explanations.
- When reading files, request specific line ranges, not entire files.

## TOOLS (30)

**File I/O**: read_file(path,start_line?,end_line?), create_file(path,content), write_file(path,content), patch_file(path,old_str,new_str), insert_lines(path,after_line,content), delete_lines(path,start_line,end_line), delete_file(path), move_file(source,destination)
**Compound**: scaffold_module(module_name,prefix,endpoints,imports?), create_page(page_name,route_path,title,api_endpoints?,content?,icon?,nav_section?)
**DB**: get_schema(collection), run_query(query_type)
**Infra**: restart_service(service), install_package(package,manager?), check_logs(service,lines?), run_tests(test_path?)
**Verification**: verify_deployment(checks[]), test_api(method,url,body?)
**Search**: grep_search(pattern,directory?,file_ext?), list_files(directory), run_command(command)
**Research**: web_search(query,max_results?), crawl_url(url), take_screenshot(url,full_page?,wait_ms?)
**Config**: manage_env(action,file?,key?,value?)
**Quality**: lint_code(path,fix?)
**Git**: git_info(action,file?)
**Subagents**: call_subagent(agent_type,task,context?) — types: tester, designer, integrator, troubleshooter
**Batch**: batch_operations(operations[]) — parallel file ops: create/write/delete/move/patch/read, max 20
**Image**: generate_image(prompt,size?)

## TOOL CALL FORMAT
```TOOL_CALL
{"tool": "tool_name", "args": {...}}
```
Multiple tool calls in ONE response = PARALLEL execution.

```DONE
Summary of what was accomplished
```
```QUESTION
Clarifying question (ONLY when genuinely blocked)
```

## CODE PATTERNS
- Routes: `router = APIRouter(prefix="/x")` + `set_db(database)` — NEVER create own motor client
- IDs: `str(uuid.uuid4())` | Timestamps: `datetime.now(timezone.utc).isoformat()`
- ALWAYS exclude `_id`: `{"_id": 0}` in projection
- Frontend: `import { API } from '../App'` then `fetch(\`\${API}/endpoint\`)` — API already includes /api
- Lucide React icons, Shadcn/UI, data-testid on all elements

## EXECUTION RULES
1. **ACT IMMEDIATELY** — Start executing tool calls. Never ask "shall I proceed?"
2. **Batch aggressively** — Multiple independent tool calls in ONE response.
3. **Self-heal on failure** — Read logs, diagnose, fix. Don't stop and report errors.
4. **ALWAYS verify** — End with verify_deployment or test_api before DONE.
5. **Maximum """ + str(MAX_ITERATIONS) + """ iterations.**

GST: intra=CGST+SGST, inter=IGST. Export=zero-rated LUT. TDS: 194J(10%), 194C(2%), 194I(10%). Revenue Ind AS 115: FP=POC, T&M=right-to-invoice, Milestone=acceptance, Retainer=straight-line."""

BA_ONLY_SUFFIX = "\n\nMODE: Business Analysis Only. Focus on requirements, compliance, accounting. No code generation."
DEV_ONLY_SUFFIX = "\n\nMODE: Coding Only. Read files, generate code, deploy. Skip business analysis."
QA_ONLY_SUFFIX = "\n\nMODE: Testing/Validation Only. Run queries, test APIs, check data integrity."

# ══════════════════════════════════════════════════════════
# TOOL EXECUTION ENGINE
# ══════════════════════════════════════════════════════════

WRITE_TOOLS = {"write_file", "create_file", "patch_file", "insert_lines", "delete_lines", "scaffold_module", "create_page", "delete_file", "move_file", "manage_env", "batch_operations"}
READ_TOOLS = {"read_file", "grep_search", "list_files", "run_command", "get_schema", "check_logs", "run_query", "verify_deployment", "web_search", "take_screenshot", "lint_code", "crawl_url", "git_info", "call_subagent", "generate_image"}

async def execute_tool(tool_name, args):
    try:
        if tool_name == "read_file":
            path = args.get("path", "")
            start_line = args.get("start_line", 1)
            end_line = args.get("end_line")
            if not is_safe_path(path):
                return {"status": "error", "error": "Access denied — blocked path"}
            if not os.path.isfile(path):
                return {"status": "error", "error": f"File not found: {path}"}
            with open(path, "r") as f:
                lines = f.readlines()
            total = len(lines)
            s = max(1, start_line) - 1
            e = min(end_line or total, total)
            numbered = [f"{i+s+1}| {line}" for i, line in enumerate(lines[s:e])]
            content = "".join(numbered)
            if len(content) > 30000:
                content = content[:30000] + "\n... [TRUNCATED] ..."
            return {"status": "ok", "path": path, "total_lines": total, "showing": f"{s+1}-{e}", "content": content}

        elif tool_name == "create_file":
            path = args.get("path", "")
            content = args.get("content", "")
            if not is_safe_path(path):
                return {"status": "error", "error": "Access denied — blocked path"}
            if os.path.isfile(path):
                return {"status": "error", "error": f"File already exists: {path}. Use patch_file."}
            os.makedirs(os.path.dirname(path), exist_ok=True)
            with open(path, "w") as f:
                f.write(content)
            await _audit_file_write(path, content, "CREATE")
            return {"status": "ok", "path": path, "size": len(content)}

        elif tool_name == "patch_file":
            path = args.get("path", "")
            old_str = args.get("old_str", "")
            new_str = args.get("new_str", "")
            if not is_safe_path(path):
                return {"status": "error", "error": "Access denied — blocked path"}
            if not os.path.isfile(path):
                return {"status": "error", "error": f"File not found: {path}"}
            with open(path, "r") as f:
                content = f.read()
            if old_str not in content:
                stripped_old = old_str.strip()
                found = any(stripped_old in line.strip() for line in content.split("\n"))
                if not found:
                    return {"status": "error", "error": "old_str not found. Use read_file first.", "hint": "Read the file to get exact text."}
            new_content = content.replace(old_str, new_str, 1)
            with open(path, "w") as f:
                f.write(new_content)
            await _audit_file_write(path, f"PATCH: {len(old_str)}→{len(new_str)} chars", "PATCH")
            return {"status": "ok", "path": path, "chars_removed": len(old_str), "chars_added": len(new_str)}

        elif tool_name == "insert_lines":
            path = args.get("path", "")
            after_line = args.get("after_line", 0)
            content = args.get("content", "")
            if not is_safe_path(path):
                return {"status": "error", "error": "Access denied — blocked path"}
            if not os.path.isfile(path):
                return {"status": "error", "error": f"File not found: {path}"}
            with open(path, "r") as f:
                lines = f.readlines()
            insert_pos = min(max(0, after_line), len(lines))
            new_lines = content.split("\n")
            for i, nl in enumerate(new_lines):
                lines.insert(insert_pos + i, nl + "\n")
            with open(path, "w") as f:
                f.writelines(lines)
            await _audit_file_write(path, f"INSERT: {len(new_lines)} lines after L{after_line}", "INSERT")
            return {"status": "ok", "path": path, "lines_inserted": len(new_lines), "at_line": after_line + 1}

        elif tool_name == "delete_lines":
            path = args.get("path", "")
            start_line = args.get("start_line", 1)
            end_line = args.get("end_line", 1)
            if not is_safe_path(path):
                return {"status": "error", "error": "Access denied — blocked path"}
            if not os.path.isfile(path):
                return {"status": "error", "error": f"File not found: {path}"}
            with open(path, "r") as f:
                lines = f.readlines()
            s = max(1, start_line) - 1
            e = min(end_line, len(lines))
            new_lines = lines[:s] + lines[e:]
            with open(path, "w") as f:
                f.writelines(new_lines)
            await _audit_file_write(path, f"DELETE: lines {start_line}-{end_line}", "DELETE")
            return {"status": "ok", "path": path, "lines_deleted": e - s}

        elif tool_name == "write_file":
            path = args.get("path", "")
            content = args.get("content", "")
            if not is_safe_path(path):
                return {"status": "error", "error": "Access denied — blocked path"}
            if os.path.isfile(path):
                with open(path, "r") as f:
                    existing = f.readlines()
                if len(existing) > 50:
                    return {"status": "error", "error": f"File has {len(existing)} lines. Use patch_file."}
            os.makedirs(os.path.dirname(path), exist_ok=True)
            with open(path, "w") as f:
                f.write(content)
            await _audit_file_write(path, content, "WRITE")
            return {"status": "ok", "path": path, "size": len(content)}

        # ── COMPOUND TOOLS ──

        elif tool_name == "scaffold_module":
            return await _scaffold_module(args)

        elif tool_name == "create_page":
            return await _create_page(args)

        # ── DATABASE ──

        elif tool_name == "run_query":
            """Execute MongoDB queries directly. Supports find, count, insert, update, delete, aggregate, and full health check."""
            query_type = args.get("query_type", args.get("operation", "full_health_check"))
            collection_name = args.get("collection", "")

            if query_type == "full_health_check":
                result = await _run_test_query("full_health_check")
                return {"status": "ok", "query_type": "full_health_check", "results": result}

            if not collection_name:
                return {"status": "error", "error": "collection name required"}

            coll = db[collection_name]
            query = args.get("query", {})
            try:
                if query_type in ["find", "read"]:
                    projection = args.get("projection", {"_id": 0})
                    if "_id" not in projection:
                        projection["_id"] = 0
                    limit = min(args.get("limit", 20), 100)
                    docs = await coll.find(query, projection).limit(limit).to_list(limit)
                    return {"status": "ok", "collection": collection_name, "count": len(docs), "documents": docs}

                elif query_type == "count":
                    count = await coll.count_documents(query)
                    return {"status": "ok", "collection": collection_name, "count": count}

                elif query_type in ["insert", "insert_one"]:
                    doc = args.get("document", args.get("doc", {}))
                    if not doc:
                        return {"status": "error", "error": "document required for insert"}
                    await coll.insert_one(doc)
                    doc.pop("_id", None)
                    return {"status": "ok", "collection": collection_name, "inserted": doc}

                elif query_type in ["insert_many"]:
                    docs = args.get("documents", [])
                    if not docs:
                        return {"status": "error", "error": "documents array required"}
                    await coll.insert_many(docs)
                    for d in docs:
                        d.pop("_id", None)
                    return {"status": "ok", "collection": collection_name, "inserted_count": len(docs)}

                elif query_type in ["update", "update_many"]:
                    update = args.get("update", {})
                    if not update:
                        return {"status": "error", "error": "update object required"}
                    result = await coll.update_many(query, update)
                    return {"status": "ok", "collection": collection_name, "matched": result.matched_count, "modified": result.modified_count}

                elif query_type in ["update_one"]:
                    update = args.get("update", {})
                    result = await coll.update_one(query, update)
                    return {"status": "ok", "collection": collection_name, "matched": result.matched_count, "modified": result.modified_count}

                elif query_type in ["delete", "delete_many"]:
                    result = await coll.delete_many(query)
                    return {"status": "ok", "collection": collection_name, "deleted": result.deleted_count}

                elif query_type in ["delete_one"]:
                    result = await coll.delete_one(query)
                    return {"status": "ok", "collection": collection_name, "deleted": result.deleted_count}

                elif query_type == "aggregate":
                    pipeline = args.get("pipeline", [])
                    if not pipeline:
                        return {"status": "error", "error": "pipeline array required"}
                    docs = await coll.aggregate(pipeline).to_list(100)
                    for d in docs:
                        d.pop("_id", None)
                    return {"status": "ok", "collection": collection_name, "results": docs, "count": len(docs)}

                elif query_type == "distinct":
                    field = args.get("field", "")
                    if not field:
                        return {"status": "error", "error": "field required for distinct"}
                    values = await coll.distinct(field, query)
                    return {"status": "ok", "collection": collection_name, "field": field, "values": values, "count": len(values)}

                elif query_type == "drop":
                    await coll.drop()
                    return {"status": "ok", "collection": collection_name, "dropped": True}

                else:
                    return {"status": "error", "error": f"Unknown query_type: {query_type}. Use: find, count, insert, insert_many, update, update_one, update_many, delete, delete_one, delete_many, aggregate, distinct, drop"}
            except Exception as e:
                return {"status": "error", "error": str(e)}

        elif tool_name == "get_schema":
            collection = args.get("collection", "")
            if not collection:
                return {"status": "error", "error": "collection name required"}
            try:
                sample = await db[collection].find_one({}, {"_id": 0})
                if not sample:
                    return {"status": "ok", "collection": collection, "fields": [], "note": "Collection empty"}
                fields = {k: type(v).__name__ for k, v in sample.items()}
                count = await db[collection].count_documents({})
                return {"status": "ok", "collection": collection, "count": count, "fields": fields, "sample_keys": list(sample.keys())}
            except Exception as ex:
                return {"status": "error", "error": str(ex)}

        # ── INFRASTRUCTURE ──

        elif tool_name == "restart_service":
            service = args.get("service", "backend")
            if service not in ["backend", "frontend"]:
                return {"status": "error", "error": "Can only restart 'backend' or 'frontend'"}
            proc = subprocess.run(["sudo", "supervisorctl", "restart", service], capture_output=True, text=True, timeout=15)
            wait_time = 4 if service == "backend" else 8  # Frontend rebuild takes longer
            await asyncio.sleep(wait_time)
            return {"status": "ok", "service": service, "output": proc.stdout.strip()}

        elif tool_name == "test_api":
            method = args.get("method", "GET").upper()
            url_path = args.get("url", "")
            body = args.get("body")
            base_url = "http://localhost:8001"
            full_url = f"{base_url}{url_path}"
            async with httpx.AsyncClient(timeout=15) as client:
                if method == "GET":
                    resp = await client.get(full_url)
                elif method == "POST":
                    resp = await client.post(full_url, json=body)
                elif method == "PUT":
                    resp = await client.put(full_url, json=body)
                elif method == "DELETE":
                    resp = await client.delete(full_url)
                else:
                    return {"status": "error", "error": f"Unsupported method: {method}"}
            resp_body = resp.text[:3000]
            try:
                resp_body = resp.json()
                if isinstance(resp_body, list) and len(resp_body) > 5:
                    resp_body = {"count": len(resp_body), "sample": resp_body[:3], "note": f"...{len(resp_body)} total"}
            except Exception:
                pass
            return {"status": "ok", "http_status": resp.status_code, "url": url_path, "response": resp_body}

        elif tool_name == "check_logs":
            service = args.get("service", "backend")
            lines = args.get("lines", 50)
            log_files = {
                "backend": "/var/log/supervisor/backend.err.log",
                "frontend": "/var/log/supervisor/frontend.err.log",
                "backend_out": "/var/log/supervisor/backend.out.log",
                "frontend_out": "/var/log/supervisor/frontend.out.log",
            }
            log_path = log_files.get(service, log_files["backend"])
            if not os.path.isfile(log_path):
                return {"status": "error", "error": f"Log file not found: {log_path}"}
            try:
                proc = subprocess.run(["tail", "-n", str(min(lines, 200)), log_path], capture_output=True, text=True, timeout=5)
                return {"status": "ok", "service": service, "lines": proc.stdout[-8000:] if proc.stdout else "(empty)"}
            except Exception as e:
                return {"status": "error", "error": str(e)}

        elif tool_name == "install_package":
            package = args.get("package", "")
            manager = args.get("manager", "pip")
            if not package or not re.match(r'^[a-zA-Z0-9\-_.=<>!@\[\],\s]+$', package):
                return {"status": "error", "error": "Invalid package name"}
            try:
                pkg_list = shlex.split(package)
                if manager == "pip":
                    proc = subprocess.run(["pip", "install"] + pkg_list, capture_output=True, text=True, timeout=60, cwd="/app/backend")
                    if proc.returncode == 0:
                        subprocess.run(["sh", "-c", "pip freeze > /app/backend/requirements.txt"], timeout=10)
                    return {"status": "ok" if proc.returncode == 0 else "error", "package": package, "output": proc.stdout[-2000:]}
                elif manager == "yarn":
                    proc = subprocess.run(["yarn", "add"] + pkg_list, capture_output=True, text=True, timeout=90, cwd="/app/frontend")
                    return {"status": "ok" if proc.returncode == 0 else "error", "package": package, "output": proc.stdout[-2000:]}
                else:
                    return {"status": "error", "error": f"Unknown manager: {manager}"}
            except subprocess.TimeoutExpired:
                return {"status": "error", "error": "Installation timed out"}

        elif tool_name == "run_tests":
            test_path = args.get("test_path", "/app/backend/tests/")
            if not re.match(r'^[a-zA-Z0-9/_.\-]+$', test_path):
                return {"status": "error", "error": "Invalid test path"}
            try:
                proc = subprocess.run(
                    ["python", "-m", "pytest", test_path, "-v", "--tb=short", "--no-header", "-q"],
                    capture_output=True, text=True, timeout=60, cwd="/app")
                output = proc.stdout[-5000:]
                if proc.stderr:
                    output += proc.stderr[-1000:]
                return {"status": "ok", "output": output, "exit_code": proc.returncode}
            except subprocess.TimeoutExpired:
                return {"status": "error", "error": "Tests timed out"}

        # ── SEARCH ──

        elif tool_name == "grep_search":
            pattern = args.get("pattern", "")
            directory = args.get("directory", "/app/backend")
            file_ext = args.get("file_ext", "")
            if not pattern:
                return {"status": "error", "error": "pattern is required"}
            if not is_safe_path(directory):
                return {"status": "error", "error": "Access denied"}
            grep_args = ["grep", "-rn", "-E", "-i"]
            if file_ext:
                grep_args.extend(["--include", f"*.{file_ext}"])
            else:
                for ext in ["py", "js", "jsx", "ts", "tsx", "css", "json"]:
                    grep_args.extend(["--include", f"*.{ext}"])
            grep_args.extend([pattern, directory])
            try:
                proc = subprocess.run(grep_args, capture_output=True, text=True, timeout=10, cwd="/app")
                raw_lines = proc.stdout.strip().split("\n") if proc.stdout.strip() else []
                matches = []
                for line in raw_lines[:60]:
                    parts = line.split(":", 2)
                    if len(parts) >= 3:
                        matches.append({"file": parts[0], "line": parts[1], "text": parts[2].strip()[:200]})
                    else:
                        matches.append({"text": line[:200]})
                return {"status": "ok", "pattern": pattern, "matches": matches, "count": len(matches)}
            except subprocess.TimeoutExpired:
                return {"status": "error", "error": "Search timed out"}

        elif tool_name == "list_files":
            directory = args.get("directory", "/app/backend")
            if not is_safe_path(directory):
                return {"status": "error", "error": "Access denied"}
            files = []
            for f in sorted(glob.glob(f"{directory}/**", recursive=True)):
                if os.path.isfile(f) and is_safe_path(f):
                    ext = os.path.splitext(f)[1]
                    if ext in [".py", ".js", ".jsx", ".ts", ".tsx", ".css", ".json", ".md"]:
                        files.append({"path": f, "relative": f.replace("/app/", ""), "size": os.path.getsize(f)})
            return {"status": "ok", "files": files[:100], "count": len(files)}

        elif tool_name == "run_command":
            cmd = args.get("command", "")
            timeout_secs = min(args.get("timeout", 60), 120)
            HARD_BLOCKED = ["rm -rf /", "mkfs", ":(){", "dd if=", "curl|sh", "wget|sh", "curl|bash", "wget|bash"]
            for bc in HARD_BLOCKED:
                if bc in cmd:
                    return {"status": "error", "error": f"Command blocked for safety: contains '{bc}'"}
            try:
                proc = subprocess.run(
                    ["bash", "-c", cmd],
                    capture_output=True, text=True, timeout=timeout_secs, cwd="/app",
                    env={**os.environ, "PATH": os.environ.get("PATH", "/usr/bin:/bin")})
                output = proc.stdout[:8000]
                if proc.stderr:
                    output += f"\n[STDERR]: {proc.stderr[:2000]}"
                return {"status": "ok", "command": cmd, "output": output, "exit_code": proc.returncode}
            except subprocess.TimeoutExpired:
                return {"status": "error", "error": f"Command timed out ({timeout_secs}s)"}

        elif tool_name == "verify_deployment":
            """Comprehensive deployment verification — checks backend health, specific API endpoints, and frontend routes."""
            checks = args.get("checks", [])
            if not checks:
                checks = [{"type": "backend_health"}]
            results = []
            for check in checks[:8]:
                check_type = check.get("type", "api")
                if check_type == "backend_health":
                    # Check if backend is running and responding
                    try:
                        log_proc = subprocess.run(["tail", "-n", "10", "/var/log/supervisor/backend.err.log"], capture_output=True, text=True, timeout=5)
                        startup_ok = "Application startup complete" in (log_proc.stdout or "")
                        async with httpx.AsyncClient(timeout=10) as client:
                            resp = await client.get("http://localhost:8001/api/health")
                            health_ok = resp.status_code == 200
                        results.append({"check": "backend_health", "startup_ok": startup_ok, "health_endpoint": health_ok, "status": "pass" if (startup_ok and health_ok) else "fail"})
                    except Exception as e:
                        results.append({"check": "backend_health", "status": "fail", "error": str(e)})
                elif check_type == "api":
                    url_path = check.get("url", "")
                    method = check.get("method", "GET").upper()
                    expected_status = check.get("expected_status", 200)
                    try:
                        base_url = "http://localhost:8001"
                        full_url = f"{base_url}{url_path}"
                        async with httpx.AsyncClient(timeout=10) as client:
                            if method == "GET":
                                resp = await client.get(full_url)
                            elif method == "POST":
                                resp = await client.post(full_url, json=check.get("body", {}))
                            else:
                                resp = await client.request(method, full_url)
                        passed = resp.status_code == expected_status
                        body_preview = resp.text[:500]
                        try:
                            body_preview = resp.json()
                            if isinstance(body_preview, list):
                                body_preview = {"count": len(body_preview), "sample": body_preview[:2]}
                        except Exception:
                            pass
                        results.append({"check": "api", "url": url_path, "method": method, "http_status": resp.status_code, "expected": expected_status, "status": "pass" if passed else "fail", "response_preview": body_preview})
                    except Exception as e:
                        results.append({"check": "api", "url": url_path, "status": "fail", "error": str(e)})
                elif check_type == "frontend_route":
                    route = check.get("route", "")
                    try:
                        # Check if the route is registered in App.js
                        with open("/app/frontend/src/App.js", "r") as f:
                            app_content = f.read()
                        route_exists = f'path="{route}"' in app_content
                        results.append({"check": "frontend_route", "route": route, "registered": route_exists, "status": "pass" if route_exists else "fail"})
                    except Exception as e:
                        results.append({"check": "frontend_route", "route": route, "status": "fail", "error": str(e)})
                elif check_type == "file_exists":
                    path = check.get("path", "")
                    exists = os.path.isfile(path)
                    size = os.path.getsize(path) if exists else 0
                    results.append({"check": "file_exists", "path": path, "exists": exists, "size": size, "status": "pass" if exists else "fail"})
            all_passed = all(r.get("status") == "pass" for r in results)
            return {"status": "ok", "all_passed": all_passed, "checks": results, "summary": f"{sum(1 for r in results if r['status']=='pass')}/{len(results)} checks passed"}

        elif tool_name == "web_search":
            """Search the web using DuckDuckGo. Returns top results with titles, URLs, and snippets."""
            query = args.get("query", "")
            max_results = min(args.get("max_results", 5), 10)
            if not query:
                return {"status": "error", "error": "query is required"}
            try:
                from ddgs import DDGS
                raw_results = list(DDGS().text(query, max_results=max_results))
                results = []
                for r in raw_results:
                    results.append({
                        "title": r.get("title", ""),
                        "url": r.get("href", ""),
                        "snippet": r.get("body", "")[:400],
                    })
                return {"status": "ok", "query": query, "results": results, "count": len(results)}
            except Exception as e:
                return {"status": "error", "error": f"Web search failed: {str(e)}"}

        elif tool_name == "take_screenshot":
            """Take a screenshot of a URL using Playwright. Returns the image path and a base64 thumbnail."""
            url = args.get("url", "")
            full_page = args.get("full_page", False)
            wait_ms = min(args.get("wait_ms", 2000), 10000)
            if not url:
                return {"status": "error", "error": "url is required"}
            # Default to local frontend if path-only
            if url.startswith("/"):
                url = f"http://localhost:3000{url}"
            elif not url.startswith("http"):
                url = f"http://localhost:3000/{url}"
            screenshot_id = str(uuid.uuid4())[:8]
            screenshot_path = f"/app/backend/uploads/screenshot_{screenshot_id}.png"
            try:
                proc = await asyncio.create_subprocess_exec(
                    "python3", "/app/backend/screenshot_helper.py",
                    url, screenshot_path, str(full_page), str(wait_ms),
                    stdout=asyncio.subprocess.PIPE,
                    stderr=asyncio.subprocess.PIPE,
                )
                stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=20)
                if proc.returncode != 0:
                    err = stderr.decode()[-500:] if stderr else "Unknown error"
                    return {"status": "error", "error": f"Screenshot failed: {err}"}
                # Generate small base64 preview (first 200KB of file)
                file_size = os.path.getsize(screenshot_path) if os.path.isfile(screenshot_path) else 0
                return {
                    "status": "ok",
                    "path": screenshot_path,
                    "url_captured": url,
                    "file_size_kb": round(file_size / 1024, 1),
                    "full_page": full_page,
                    "note": "Screenshot saved. Use read_file or serve from /uploads/ to view.",
                }
            except asyncio.TimeoutError:
                try:
                    proc.kill()
                except Exception:
                    pass
                return {"status": "error", "error": "Screenshot timed out (20s limit). Try a simpler URL or increase wait_ms."}
            except Exception as e:
                return {"status": "error", "error": f"Screenshot error: {str(e)}"}

        elif tool_name == "delete_file":
            """Delete a file from the project."""
            path = args.get("path", "")
            if not is_safe_path(path):
                return {"status": "error", "error": "Access denied — blocked path"}
            if not os.path.isfile(path):
                return {"status": "error", "error": f"File not found: {path}"}
            try:
                size = os.path.getsize(path)
                os.remove(path)
                await _audit_file_write(path, f"DELETED ({size} bytes)", "DELETE_FILE")
                return {"status": "ok", "path": path, "deleted": True, "size_was": size}
            except Exception as e:
                return {"status": "error", "error": str(e)}

        elif tool_name == "move_file":
            """Move or rename a file."""
            source = args.get("source", "")
            destination = args.get("destination", "")
            if not is_safe_path(source) or not is_safe_path(destination):
                return {"status": "error", "error": "Access denied — blocked path"}
            if not os.path.isfile(source):
                return {"status": "error", "error": f"Source not found: {source}"}
            try:
                import shutil
                os.makedirs(os.path.dirname(destination), exist_ok=True)
                shutil.move(source, destination)
                await _audit_file_write(destination, f"MOVED from {source}", "MOVE_FILE")
                return {"status": "ok", "source": source, "destination": destination}
            except Exception as e:
                return {"status": "error", "error": str(e)}

        elif tool_name == "manage_env":
            """Read or write .env file variables safely."""
            action = args.get("action", "read")  # read, set, delete
            env_file = args.get("file", "backend")  # backend or frontend
            env_path = "/app/backend/.env" if env_file == "backend" else "/app/frontend/.env"

            if action == "read":
                if not os.path.isfile(env_path):
                    return {"status": "error", "error": f".env not found: {env_path}"}
                with open(env_path, "r") as f:
                    lines = f.readlines()
                # Mask sensitive values
                vars_list = []
                for line in lines:
                    line = line.strip()
                    if "=" in line and not line.startswith("#"):
                        key = line.split("=", 1)[0]
                        val = line.split("=", 1)[1]
                        masked = val[:4] + "..." + val[-4:] if len(val) > 12 else val
                        vars_list.append({"key": key, "value_preview": masked, "length": len(val)})
                return {"status": "ok", "file": env_path, "variables": vars_list}

            elif action == "set":
                key = args.get("key", "")
                value = args.get("value", "")
                if not key or not re.match(r'^[A-Z_][A-Z0-9_]*$', key):
                    return {"status": "error", "error": "Invalid env key. Must be UPPER_SNAKE_CASE."}
                # Protected keys that cannot be overwritten
                PROTECTED = {"MONGO_URL", "DB_NAME", "REACT_APP_BACKEND_URL"}
                if key in PROTECTED:
                    return {"status": "error", "error": f"Cannot modify protected key: {key}"}
                # Read existing
                lines = []
                if os.path.isfile(env_path):
                    with open(env_path, "r") as f:
                        lines = f.readlines()
                # Update or append
                found = False
                for i, line in enumerate(lines):
                    if line.strip().startswith(f"{key}="):
                        lines[i] = f"{key}={value}\n"
                        found = True
                        break
                if not found:
                    lines.append(f"{key}={value}\n")
                with open(env_path, "w") as f:
                    f.writelines(lines)
                return {"status": "ok", "file": env_path, "key": key, "action": "updated" if found else "added"}

            elif action == "delete":
                key = args.get("key", "")
                PROTECTED = {"MONGO_URL", "DB_NAME", "REACT_APP_BACKEND_URL", "EMERGENT_LLM_KEY"}
                if key in PROTECTED:
                    return {"status": "error", "error": f"Cannot delete protected key: {key}"}
                if os.path.isfile(env_path):
                    with open(env_path, "r") as f:
                        lines = f.readlines()
                    new_lines = [line for line in lines if not line.strip().startswith(f"{key}=")]
                    with open(env_path, "w") as f:
                        f.writelines(new_lines)
                return {"status": "ok", "file": env_path, "key": key, "action": "deleted"}
            else:
                return {"status": "error", "error": f"Unknown action: {action}. Use read, set, or delete."}

        elif tool_name == "lint_code":
            """Run linters on code files. Uses ruff for Python, eslint for JS/JSX/TS/TSX."""
            path = args.get("path", "")
            fix = args.get("fix", False)
            if not path:
                return {"status": "error", "error": "path is required"}
            if not re.match(r'^[a-zA-Z0-9/_.\-]+$', path):
                return {"status": "error", "error": "Invalid path characters"}
            ext = os.path.splitext(path)[1].lower()
            try:
                if ext == ".py" or (os.path.isdir(path) and not path.endswith("src")):
                    cmd = ["ruff", "check", path]
                    if fix:
                        cmd.append("--fix")
                    proc = subprocess.run(cmd, capture_output=True, text=True, timeout=30, cwd="/app")
                    issues = proc.stdout.strip() if proc.stdout.strip() else "No issues found"
                    return {"status": "ok", "linter": "ruff", "path": path, "output": issues[:3000], "exit_code": proc.returncode}
                elif ext in [".js", ".jsx", ".ts", ".tsx"] or path.endswith("src"):
                    cmd = ["npx", "eslint", path]
                    if fix:
                        cmd.append("--fix")
                    proc = subprocess.run(cmd, capture_output=True, text=True, timeout=30, cwd="/app/frontend")
                    issues = proc.stdout.strip() if proc.stdout.strip() else "No issues found"
                    return {"status": "ok", "linter": "eslint", "path": path, "output": issues[:3000], "exit_code": proc.returncode}
                else:
                    return {"status": "error", "error": f"No linter for extension: {ext}"}
            except subprocess.TimeoutExpired:
                return {"status": "error", "error": "Linting timed out"}

        elif tool_name == "crawl_url":
            """Fetch and extract text content from a URL for research."""
            url = args.get("url", "")
            if not url or not url.startswith("http"):
                return {"status": "error", "error": "Valid HTTP URL required"}
            try:
                async with httpx.AsyncClient(timeout=15, follow_redirects=True) as client:
                    resp = await client.get(url, headers={"User-Agent": "Mozilla/5.0 Kairos Engine"})
                text = resp.text
                # Strip HTML tags for cleaner text
                import re as re_mod
                clean = re_mod.sub(r'<script[^>]*>.*?</script>', '', text, flags=re_mod.DOTALL)
                clean = re_mod.sub(r'<style[^>]*>.*?</style>', '', clean, flags=re_mod.DOTALL)
                clean = re_mod.sub(r'<[^>]+>', ' ', clean)
                clean = re_mod.sub(r'\s+', ' ', clean).strip()
                return {"status": "ok", "url": url, "http_status": resp.status_code, "content": clean[:8000], "full_length": len(clean)}
            except Exception as e:
                return {"status": "error", "error": f"Crawl failed: {str(e)}"}

        elif tool_name == "git_info":
            """Get git status, recent commits, or diff information."""
            action = args.get("action", "log")  # log, status, diff
            try:
                if action == "log":
                    proc = subprocess.run(["git", "log", "--oneline", "-20"], capture_output=True, text=True, timeout=10, cwd="/app")
                elif action == "status":
                    proc = subprocess.run(["git", "status", "--short"], capture_output=True, text=True, timeout=10, cwd="/app")
                elif action == "diff":
                    file_path = args.get("file", "")
                    git_cmd = ["git", "diff", "HEAD"]
                    if file_path:
                        git_cmd.extend(["--", file_path])
                    else:
                        git_cmd.append("--stat")
                    proc = subprocess.run(git_cmd, capture_output=True, text=True, timeout=10, cwd="/app")
                else:
                    return {"status": "error", "error": f"Unknown git action: {action}"}
                return {"status": "ok", "action": action, "output": proc.stdout[:5000]}
            except Exception as e:
                return {"status": "error", "error": str(e)}

        elif tool_name == "call_subagent":
            """Call a specialized AI subagent for testing, design, integration, or troubleshooting."""
            agent_type = args.get("agent_type", "")
            task = args.get("task", "")
            context = args.get("context", "")
            if not agent_type or not task:
                return {"status": "error", "error": "agent_type and task are required"}
            return await call_subagent(agent_type, task, context)

        elif tool_name == "batch_operations":
            """Execute multiple file operations in parallel. Each op: {action, path, content?, destination?}."""
            operations = args.get("operations", [])
            if not operations:
                return {"status": "error", "error": "operations array required"}

            async def _do_op(op):
                action = op.get("action", "")
                path = op.get("path", "")
                try:
                    if action == "create":
                        return await execute_tool("create_file", {"path": path, "content": op.get("content", "")})
                    elif action == "write":
                        return await execute_tool("write_file", {"path": path, "content": op.get("content", "")})
                    elif action == "delete":
                        return await execute_tool("delete_file", {"path": path})
                    elif action == "move":
                        return await execute_tool("move_file", {"path": path, "source": path, "destination": op.get("destination", "")})
                    elif action == "patch":
                        return await execute_tool("patch_file", {"path": path, "search": op.get("search", ""), "replace": op.get("replace", "")})
                    elif action == "read":
                        return await execute_tool("read_file", {"path": path})
                    else:
                        return {"status": "error", "error": f"Unknown action: {action}"}
                except Exception as e:
                    return {"status": "error", "path": path, "error": str(e)}

            results = await asyncio.gather(*[_do_op(op) for op in operations[:20]])
            succeeded = sum(1 for r in results if r.get("status") == "ok")
            return {
                "status": "ok",
                "total": len(operations),
                "succeeded": succeeded,
                "failed": len(operations) - succeeded,
                "results": [{"action": op.get("action"), "path": op.get("path"), "status": r.get("status"), "error": r.get("error")} for op, r in zip(operations, results)],
            }

        elif tool_name == "generate_image":
            """Generate an image from a text prompt using AI (GPT Image 1)."""
            prompt = args.get("prompt", "")
            size = args.get("size", "1024x1024")
            if not prompt:
                return {"status": "error", "error": "prompt is required"}
            return await gen_image(prompt, size)

        else:
            return {"status": "error", "error": f"Unknown tool: {tool_name}"}
    except Exception as e:
        return {"status": "error", "error": str(e)}


# ══════════════════════════════════════════════════════════
# COMPOUND TOOLS
# ══════════════════════════════════════════════════════════

async def _scaffold_module(args):
    """Create a complete backend module: route file + server.py registration + restart + verify + auto-polish."""
    module_name = args.get("module_name", "")
    prefix = args.get("prefix", "")
    endpoints = args.get("endpoints", [])
    extra_imports = args.get("imports", "")

    if not module_name or not prefix or not endpoints:
        return {"status": "error", "error": "module_name, prefix, and endpoints are required"}

    safe_module = re.sub(r'[^a-z0-9_]', '_', module_name.lower())
    file_path = f"/app/backend/routes_{safe_module}.py"
    if os.path.isfile(file_path):
        return {"status": "error", "error": f"Module file already exists: {file_path}. Use patch_file to modify."}

    tag = module_name.replace("_", " ").title()
    code_lines = [
        f'"""Auto-generated module: {tag}"""',
        'from fastapi import APIRouter',
        'from datetime import datetime, timezone',
        'import uuid',
    ]
    if extra_imports:
        code_lines.append(extra_imports)
    code_lines.extend([
        '',
        f'router = APIRouter(prefix="{prefix}", tags=["{tag}"])',
        '',
        'db = None',
        '',
        'def set_db(database):',
        '    global db',
        '    db = database',
        '',
    ])
    for ep in endpoints:
        method = ep.get("method", "GET").lower()
        path = ep.get("path", "")
        name = ep.get("name", f"handler_{method}")
        body = ep.get("body", "    return {}")
        body_lines = body.split("\n")
        formatted_body = "\n".join(f"    {line}" if not line.startswith("    ") else line for line in body_lines)

        # Extract path parameters for proper function signatures
        import re as _re
        path_params = _re.findall(r'\{(\w+)\}', path)

        if method == "post":
            code_lines.append(f'@router.post("{path}")')
            params = ", ".join([f"{p}: str" for p in path_params] + ["body: dict"])
            code_lines.append(f'async def {name}({params}):')
        elif method == "put":
            code_lines.append(f'@router.put("{path}")')
            params = ", ".join([f"{p}: str" for p in path_params] + ["body: dict"])
            code_lines.append(f'async def {name}({params}):')
        elif method == "delete":
            ep_path = path if path else "/{item_id}"
            code_lines.append(f'@router.delete("{ep_path}")')
            del_params = path_params if path_params else ["item_id"]
            code_lines.append(f'async def {name}({", ".join(f"{p}: str" for p in del_params)}):')
        else:
            code_lines.append(f'@router.get("{path}")')
            if path_params:
                code_lines.append(f'async def {name}({", ".join(f"{p}: str" for p in path_params)}):')
            else:
                code_lines.append(f'async def {name}():')
        code_lines.append(formatted_body)
        code_lines.append('')

    file_content = "\n".join(code_lines)

    # ── AUTO-POLISH: fix known LLM code-gen bugs ──
    file_content = _polish_generated_python(file_content)

    with open(file_path, "w") as f:
        f.write(file_content)
    await _audit_file_write(file_path, file_content, "SCAFFOLD")

    # Register in server.py
    var_name = f"{safe_module}_router"
    set_fn = f"set_{safe_module}_db"
    registration_code = f"""
    # {tag}
    from routes_{safe_module} import router as {var_name}, set_db as {set_fn}
    {set_fn}(db)
    api_router.include_router({var_name})"""

    server_path = "/app/backend/server.py"
    with open(server_path, "r") as f:
        server_content = f.read()

    marker = 'logging.info("ERP modules integrated (including 10 advanced modules)")'
    if marker not in server_content:
        marker = 'logging.info("ERP modules will be integrated")'
    if marker in server_content:
        server_content = server_content.replace(marker, f"{registration_code}\n    {marker}")
        with open(server_path, "w") as f:
            f.write(server_content)
    else:
        return {"status": "partial", "path": file_path, "warning": "Could not find server.py marker. Register manually."}

    # Restart backend
    proc = subprocess.run(["sudo", "supervisorctl", "restart", "backend"], capture_output=True, text=True, timeout=15)
    await asyncio.sleep(3)

    # Verify startup
    log_proc = subprocess.run(["tail", "-n", "15", "/var/log/supervisor/backend.err.log"], capture_output=True, text=True, timeout=5)
    startup_ok = "Application startup complete" in (log_proc.stdout or "")

    # If startup failed, try auto-fix once
    auto_fix_applied = False
    if not startup_ok:
        fix_result = _auto_fix_startup_error(file_path, log_proc.stdout or "")
        if fix_result:
            auto_fix_applied = True  # noqa: F841
            subprocess.run(["sudo", "supervisorctl", "restart", "backend"], capture_output=True, text=True, timeout=15)
            await asyncio.sleep(3)
            log_proc = subprocess.run(["tail", "-n", "15", "/var/log/supervisor/backend.err.log"], capture_output=True, text=True, timeout=5)
            startup_ok = "Application startup complete" in (log_proc.stdout or "")

    # Test first GET endpoint
    test_result = None
    first_get = next((ep for ep in endpoints if ep.get("method", "GET").upper() == "GET"), None)
    if first_get and startup_ok:
        test_url = f"http://localhost:8001/api{prefix}{first_get.get('path', '')}"
        # Strip path params for test URL
        test_url = re.sub(r'/\{[^}]+\}', '', test_url)
        try:
            async with httpx.AsyncClient(timeout=10) as client:
                resp = await client.get(test_url)
                test_result = {"http_status": resp.status_code, "ok": resp.status_code < 400}
        except Exception as e:
            test_result = {"http_status": 0, "ok": False, "error": str(e)}

    return {
        "status": "ok" if startup_ok else "error",
        "path": file_path,
        "module_name": safe_module,
        "prefix": prefix,
        "endpoints_created": len(endpoints),
        "registered_in_server": True,
        "backend_restarted": True,
        "startup_ok": startup_ok,
        "auto_fix_applied": auto_fix_applied,
        "auto_polish": True,
        "test_result": test_result,
        "logs": log_proc.stdout[-500:] if not startup_ok else "",
    }


async def _create_page(args):
    """Create a React page component + register route in App.js + add sidebar nav."""
    page_name = args.get("page_name", "")
    route_path = args.get("route_path", "")
    title = args.get("title", page_name)
    api_endpoints = args.get("api_endpoints", [])
    custom_content = args.get("content", "")
    icon = args.get("icon", "FileText")
    nav_section = args.get("nav_section", "")

    if not page_name or not route_path:
        return {"status": "error", "error": "page_name and route_path required"}

    file_path = f"/app/frontend/src/pages/{page_name}.js"
    if os.path.isfile(file_path):
        return {"status": "error", "error": f"Page already exists: {file_path}. Use patch_file."}

    if custom_content:
        page_code = custom_content
    else:
        fetch_code = ""
        if api_endpoints:
            # Fix: api_endpoints should NOT include /api prefix (API const already has it)
            ep = api_endpoints[0]
            if ep.startswith("/api/"):
                ep = ep[4:]  # strip /api prefix — API const already includes it
            elif not ep.startswith("/"):
                ep = "/" + ep
            fetch_code = f"""
  const [data, setData] = useState([]);
  const [loading, setLoading] = useState(true);
  useEffect(() => {{
    fetch(`${{API}}{ep}`).then(r => r.json()).then(d => {{ setData(Array.isArray(d) ? d : []); setLoading(false); }}).catch(() => setLoading(false));
  }}, []);"""

        page_code = f"""import {{ useState, useEffect }} from 'react';
import {{ API }} from '../App';

export default function {page_name}() {{
{fetch_code}
  return (
    <div className="p-6 space-y-6" data-testid="{page_name.lower()}-page">
      <h1 className="text-2xl font-bold text-[#E8EDF2]">{title}</h1>
      {'{loading ? <p className="text-[#4A5B6E]">Loading...</p> : <pre className="text-xs text-[#c8d4e0] bg-[#0D1B2A] p-4 rounded-lg border border-[#1B2D42] overflow-auto">{JSON.stringify(data, null, 2)}</pre>}' if api_endpoints else '<p className="text-[#4A5B6E]">Content goes here</p>'}
    </div>
  );
}}
"""

    with open(file_path, "w") as f:
        f.write(page_code)

    # Register route in App.js
    app_path = "/app/frontend/src/App.js"
    with open(app_path, "r") as f:
        app_content = f.read()

    # Add import
    import_line = f"import {page_name} from './pages/{page_name}';"
    if import_line not in app_content:
        import_match = list(re.finditer(r'^import .+ from .+;$', app_content, re.MULTILINE))
        if import_match:
            last_import_end = import_match[-1].end()
            app_content = app_content[:last_import_end] + f"\n{import_line}" + app_content[last_import_end:]

    # Add route
    route_line = f'<Route path="{route_path}" element={{<{page_name} />}} />'
    if route_line not in app_content:
        routes_end = app_content.find("</Routes>")
        if routes_end != -1:
            indent = "              "
            app_content = app_content[:routes_end] + f"{indent}{route_line}\n{indent}" + app_content[routes_end:]

    # Add sidebar navigation entry
    sidebar_added = False
    if nav_section:
        # Find the nav section and add entry
        nav_marker = f"label: '{nav_section}'"
        if nav_marker not in app_content:
            # Try to add near existing nav entries
            nav_entry = f"        {{ path: '{route_path}', label: '{title}', icon: {icon} }},"
            # Find the last nav entry before a closing bracket
            last_nav = app_content.rfind("{ path: '/bank-reconciliation'")
            if last_nav == -1:
                last_nav = app_content.rfind("{ path: '/expense-management'")
            if last_nav != -1:
                line_end = app_content.find("\n", last_nav)
                if line_end != -1:
                    app_content = app_content[:line_end + 1] + nav_entry + "\n" + app_content[line_end + 1:]
                    sidebar_added = True

    with open(app_path, "w") as f:
        f.write(app_content)

    return {
        "status": "ok",
        "path": file_path,
        "page_name": page_name,
        "route_path": route_path,
        "registered_in_app_js": True,
        "sidebar_added": sidebar_added,
        "api_prefix_fixed": True,
    }


# ══════════════════════════════════════════════════════════
# AUTO-POLISH & AUTO-FIX
# ══════════════════════════════════════════════════════════

def _polish_generated_python(code: str) -> str:
    """Fix common LLM code-generation bugs in Python code targeting Motor/MongoDB."""
    original = code

    # Fix 1: .to_list() missing length argument → .to_list(500)
    code = re.sub(r'\.to_list\(\s*\)', '.to_list(500)', code)

    # Fix 2: MongoDB _id not excluded in find projections
    # Only fix if find({...}) has no projection argument
    code = re.sub(
        r'\.find\(\{([^}]*)\}\)\s*\.to_list',
        lambda m: f'.find({{{m.group(1)}}}, {{"_id": 0}}).to_list' if '"_id"' not in m.group(0) else m.group(0),
        code
    )
    # Also fix find_one without _id exclusion
    code = re.sub(
        r'\.find_one\(\{([^}]*)\}\)\s*$',
        lambda m: f'.find_one({{{m.group(1)}}}, {{"_id": 0}})' if '"_id"' not in m.group(0) else m.group(0),
        code,
        flags=re.MULTILINE,
    )

    # Fix 3: return doc that may contain _id after insert_one
    # Pattern: insert_one(var) then return var → return {k:v for k,v...}
    code = re.sub(
        r'await db\.\w+\.insert_one\((\w+)\)\s*\n(\s*)return \1\s*$',
        lambda m: f'await db.{m.group(0).split("db.")[1].split(".insert")[0]}.insert_one({m.group(1)})\n{m.group(2)}return {{k: v for k, v in {m.group(1)}.items() if k != "_id"}}',
        code,
        flags=re.MULTILINE,
    )

    # Fix 4: Missing 'body: dict' param for POST/PUT handlers
    # If function body references 'body' but signature doesn't have it
    for match in re.finditer(r'async def (\w+)\(\):\n((?:    .*\n)*)', code):
        fn_name, fn_body = match.group(1), match.group(2)
        if 'body' in fn_body and 'body' not in match.group(0).split('(')[1]:
            code = code.replace(f'async def {fn_name}():', f'async def {fn_name}(body: dict):', 1)

    # Fix 5: Aggregation _id field handling — rename _id to meaningful name
    # $group results have _id as the group key, not a MongoDB ObjectId
    code = re.sub(
        r'\{k:\s*v\s+for\s+k,\s*v\s+in\s+item\.items\(\)\s+if\s+k\s*!=\s*"_id"\}',
        '{**{("name" if k == "_id" else k): v for k, v in item.items()}}',
        code,
    )

    # Fix 6: datetime.utcnow() → datetime.now(timezone.utc)
    code = code.replace('datetime.utcnow()', 'datetime.now(timezone.utc)')

    changes = sum(1 for a, b in zip(original.split('\n'), code.split('\n')) if a != b)
    if changes > 0:
        logging.info(f"Auto-polish: fixed {changes} lines in generated code")

    return code


def _auto_fix_startup_error(file_path: str, log_output: str) -> bool:
    """Try to auto-fix common startup errors. Returns True if a fix was applied."""
    if not os.path.isfile(file_path):
        return False

    with open(file_path, "r") as f:
        code = f.read()

    original = code
    fixed = False

    # Fix: to_list() missing positional argument 'length'
    if "to_list() missing 1 required positional argument" in log_output:
        code = re.sub(r'\.to_list\(\s*\)', '.to_list(500)', code)
        fixed = True

    # Fix: 'body' is not defined (missing function parameter)
    if "name 'body' is not defined" in log_output:
        # Find async def without body param that uses body
        code = re.sub(
            r'(async def \w+)\(\)(:.*\n(?:    .*body.*\n))',
            r'\1(body: dict)\2',
            code,
        )
        fixed = True

    # Fix: IndentationError
    if "IndentationError" in log_output:
        # Try to fix by ensuring consistent 4-space indentation
        lines = code.split('\n')
        new_lines = []
        for line in lines:
            if line.strip() and line[0] not in (' ', '\t', '#', '@', 'd', 'f', 'i', 'r', 'a', '"', "'"):
                # Line doesn't start with expected char — might need indentation
                new_lines.append('    ' + line)
            else:
                new_lines.append(line)
        code = '\n'.join(new_lines)
        fixed = True

    # Fix: SyntaxError on specific line
    if "SyntaxError" in log_output:
        # Extract line number
        match = re.search(r'line (\d+)', log_output)
        if match:
            line_num = int(match.group(1))
            lines = code.split('\n')
            if 0 < line_num <= len(lines):
                problem_line = lines[line_num - 1]
                # Common fix: unbalanced quotes
                if problem_line.count('"') % 2 != 0:
                    lines[line_num - 1] = problem_line + '"'
                    code = '\n'.join(lines)
                    fixed = True

    if fixed and code != original:
        with open(file_path, "w") as f:
            f.write(code)
        logging.info(f"Auto-fix applied to {file_path}")
        return True

    return False


# ══════════════════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════════════════

async def _audit_file_write(path, detail, action_type):
    if db is None:
        return
    try:
        await db.audit_trail.insert_one({
            "id": str(uuid.uuid4()), "action": f"FILE_{action_type}", "module": "AI_ENGINE",
            "record_id": path, "record_name": os.path.basename(path),
            "changes": [{"field": "content", "new_value": str(detail)[:500]}],
            "timestamp": datetime.now(timezone.utc).isoformat(), "user": "kairos-engine",
        })
    except Exception:
        pass

async def _run_test_query(query_type):
    if query_type == "tb_balance":
        coa = await db.chart_of_accounts.find({}, {"_id": 0}).to_list(100)
        dr = sum(max(0, e["opening_balance"]) for e in coa)
        cr = sum(max(0, -e["opening_balance"]) for e in coa)
        return {"total_debit": dr, "total_credit": cr, "balanced": dr == cr, "accounts": len(coa)}
    elif query_type == "entity_validation":
        v = await db.entities.find({"entity_type": "vendor"}, {"_id": 0}).to_list(100)
        c = await db.entities.find({"entity_type": "customer"}, {"_id": 0}).to_list(100)
        return {"vendors": len(v), "customers": len(c), "vendor_missing_gstin": [x["name"] for x in v if not x.get("gstin")]}
    elif query_type == "project_health":
        p = await db.projects.find({"id": {"$ne": "PRJ-INT"}}, {"_id": 0}).to_list(20)
        return {"projects": len(p)}
    elif query_type == "collection_stats":
        cols = await db.list_collection_names()
        stats = {}
        for col in sorted(cols):
            stats[col] = await db[col].count_documents({})
        return stats
    elif query_type == "full_health_check":
        coa = await db.chart_of_accounts.find({}, {"_id": 0}).to_list(100)
        dr = sum(max(0, e["opening_balance"]) for e in coa)
        cr = sum(max(0, -e["opening_balance"]) for e in coa)
        return {
            "tb_balanced": dr == cr, "tb_total": dr, "accounts": len(coa),
            "vendors": await db.entities.count_documents({"entity_type": "vendor"}),
            "customers": await db.entities.count_documents({"entity_type": "customer"}),
            "projects": await db.projects.count_documents({}),
            "employees": await db.employees.count_documents({}),
            "timesheets": await db.timesheets.count_documents({}),
            "transactions": await db.erp_transactions.count_documents({}),
        }
    else:
        return {"error": f"Unknown query: {query_type}"}


def parse_tool_calls(text):
    calls = []
    if not text:
        return calls
    parts = text.split("```TOOL_CALL")
    for part in parts[1:]:
        end = part.find("```")
        if end != -1:
            raw = part[:end].strip()
            try:
                calls.append(json.loads(raw))
            except json.JSONDecodeError:
                pass
    return calls

def parse_questions(text):
    questions = []
    if not text:
        return questions
    parts = text.split("```QUESTION")
    for part in parts[1:]:
        end = part.find("```")
        if end != -1:
            questions.append(part[:end].strip())
    return questions

def parse_done(text):
    if not text:
        return None
    if "```DONE" in text:
        parts = text.split("```DONE")
        if len(parts) > 1:
            end = parts[1].find("```")
            if end != -1:
                return parts[1][:end].strip()
            return parts[1].strip()
    return None

def _clean_response_text(text):
    if not text:
        return ""
    cleaned = re.sub(r'```TOOL_CALL[\s\S]*?```', '', text)
    cleaned = re.sub(r'```QUESTION[\s\S]*?```', '', cleaned)
    cleaned = re.sub(r'```DONE[\s\S]*?```', '', cleaned)
    cleaned = re.sub(r'\n{3,}', '\n\n', cleaned).strip()
    return cleaned


def _compress_tool_result(tool_name, result):
    """Compress tool results to reduce LLM context consumption."""
    if result.get("status") == "error":
        return {"status": "error", "error": result.get("error", "")}

    if tool_name == "read_file":
        content = result.get("content", "")
        if len(content) > 5000:
            return {"status": "ok", "path": result.get("path"), "total_lines": result.get("total_lines"),
                    "showing": result.get("showing"), "content": content[:5000] + "\n... [COMPRESSED — showing first 5000 chars]"}
        return result

    if tool_name == "list_files":
        files = result.get("files", [])
        if len(files) > 20:
            return {"status": "ok", "count": len(files), "files": files[:20], "note": f"Showing 20/{len(files)}"}
        return result

    if tool_name == "check_logs":
        log_text = result.get("lines", "")
        if len(log_text) > 3000:
            return {"status": "ok", "service": result.get("service"), "lines": log_text[-3000:], "note": "Last 3000 chars shown"}
        return result

    if tool_name == "grep_search":
        matches = result.get("matches", [])
        if len(matches) > 20:
            return {"status": "ok", "pattern": result.get("pattern"), "count": len(matches), "matches": matches[:20], "note": f"Showing 20/{len(matches)}"}
        return result

    if tool_name in ("scaffold_module", "create_page"):
        return result  # Always return full result for compound tools

    if tool_name == "run_command":
        output = result.get("output", "")
        if len(output) > 3000:
            return {"status": "ok", "command": result.get("command"), "output": output[:3000] + "\n... [COMPRESSED]", "exit_code": result.get("exit_code")}
        return result

    if tool_name == "web_search":
        # Keep web search results compact
        results_list = result.get("results", [])
        if len(results_list) > 5:
            return {"status": "ok", "query": result.get("query"), "count": len(results_list), "results": results_list[:5], "note": f"Showing 5/{len(results_list)}"}
        return result

    if tool_name == "take_screenshot":
        # Remove base64 data from compressed results — keep metadata only
        return {"status": result.get("status"), "path": result.get("path"), "url_captured": result.get("url_captured"),
                "file_size_kb": result.get("file_size_kb"), "note": result.get("note", "Screenshot captured")}

    if tool_name == "crawl_url":
        content = result.get("content", "")
        if len(content) > 4000:
            return {"status": "ok", "url": result.get("url"), "http_status": result.get("http_status"),
                    "content": content[:4000] + "\n... [TRUNCATED]", "full_length": result.get("full_length")}

    if tool_name == "run_query":
        docs = result.get("documents", result.get("results", []))
        if isinstance(docs, list) and len(docs) > 10:
            return {**result, "documents": docs[:10], "note": f"Showing 10/{len(docs)} results"}

    if tool_name == "lint_code":
        output = result.get("output", "")
        if len(output) > 3000:
            return {**result, "output": output[:3000] + "\n... [TRUNCATED]"}

    if tool_name == "call_subagent":
        response = result.get("response", "")
        if len(response) > 6000:
            return {**result, "response": response[:6000] + "\n... [TRUNCATED]", "full_length": result.get("full_length")}

    if tool_name == "batch_operations":
        results_list = result.get("results", [])
        if len(results_list) > 10:
            return {**result, "results": results_list[:10], "note": f"Showing 10/{len(results_list)} results"}

    return result


# ══════════════════════════════════════════════════════════
# FILE UPLOAD & URL CRAWLING
# ══════════════════════════════════════════════════════════

UPLOAD_DIR = "/app/backend/uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)

def _extract_pdf(path):
    import pdfplumber
    text_parts = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages[:50]):
            t = page.extract_text()
            if t:
                text_parts.append(f"--- Page {i+1} ---\n{t}")
            tables = page.extract_tables()
            for ti, table in enumerate(tables):
                text_parts.append(f"[Table {ti+1}]\n" + "\n".join([" | ".join(str(c or "") for c in row) for row in table]))
    return "\n\n".join(text_parts)

def _extract_docx(path):
    from docx import Document
    doc = Document(path)
    parts = [para.text for para in doc.paragraphs if para.text.strip()]
    for table in doc.tables:
        parts.append("[Table]\n" + "\n".join(" | ".join(cell.text for cell in row.cells) for row in table.rows))
    return "\n".join(parts)

def _extract_xlsx(path):
    from openpyxl import load_workbook
    wb = load_workbook(path, data_only=True)
    parts = []
    for name in wb.sheetnames[:10]:
        ws = wb[name]
        rows = [" | ".join(str(c or "") for c in row) for row in ws.iter_rows(max_row=200, values_only=True)]
        parts.append(f"--- Sheet: {name} ---\n" + "\n".join(rows))
    return "\n\n".join(parts)

def _extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides[:50]):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text)
            if shape.has_table:
                for row in shape.table.rows:
                    texts.append(" | ".join(cell.text for cell in row.cells))
        if texts:
            parts.append(f"--- Slide {i+1} ---\n" + "\n".join(texts))
    return "\n\n".join(parts)

def _extract_csv(path):
    import csv
    rows = []
    with open(path, "r", errors="replace") as f:
        reader = csv.reader(f)
        for i, row in enumerate(reader):
            if i > 500:
                rows.append("... [TRUNCATED]")
                break
            rows.append(" | ".join(row))
    return "\n".join(rows)

EXTRACTORS = {".pdf": _extract_pdf, ".docx": _extract_docx, ".doc": _extract_docx, ".xlsx": _extract_xlsx, ".xls": _extract_xlsx, ".pptx": _extract_pptx, ".ppt": _extract_pptx, ".csv": _extract_csv}
TEXT_EXTS = {".txt", ".md", ".json", ".xml", ".py", ".js", ".jsx", ".ts", ".tsx", ".css", ".html", ".yaml", ".yml", ".ini", ".cfg", ".log", ".sql"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".svg", ".heic", ".heif"}

@router.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    ext = os.path.splitext(file.filename or "")[1].lower()
    file_id = str(uuid.uuid4())[:8]
    safe_name = f"{file_id}_{file.filename}"
    save_path = os.path.join(UPLOAD_DIR, safe_name)
    content_bytes = await file.read()
    with open(save_path, "wb") as f:
        f.write(content_bytes)
    size_kb = len(content_bytes) / 1024
    result = {"id": file_id, "filename": file.filename, "ext": ext, "size_kb": round(size_kb, 1), "type": "unknown", "content": ""}
    try:
        if ext in EXTRACTORS:
            result["content"] = EXTRACTORS[ext](save_path)
            result["type"] = "document"
        elif ext in TEXT_EXTS:
            with open(save_path, "r", errors="replace") as f:
                result["content"] = f.read()[:50000]
            result["type"] = "text"
        elif ext in IMAGE_EXTS:
            result["type"] = "image"
            result["content"] = f"[Image: {file.filename} ({size_kb:.0f}KB)]"
        else:
            result["type"] = "binary"
            result["content"] = f"[Unsupported: {ext}]"
    except Exception as e:
        result["content"] = f"[Extraction error: {str(e)}]"
        result["type"] = "error"
    if len(result["content"]) > 40000:
        result["content"] = result["content"][:40000] + "\n... [TRUNCATED]"
    return result

@router.post("/crawl-url")
async def crawl_url(body: dict):
    url = body.get("url", "").strip()
    if not url:
        raise HTTPException(status_code=400, detail="URL is required")
    if not url.startswith(("http://", "https://")):
        url = "https://" + url
    try:
        async with httpx.AsyncClient(timeout=20, follow_redirects=True) as client:
            resp = await client.get(url, headers={"User-Agent": "Mozilla/5.0 (compatible; KairosBot/1.0)"})
            resp.raise_for_status()
        content_type = resp.headers.get("content-type", "")
        raw = resp.text
        if "html" in content_type:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(raw, "html.parser")
            for tag in soup(["script", "style", "nav", "footer", "header", "aside", "noscript"]):
                tag.decompose()
            title = soup.title.string if soup.title else url
            text = "\n".join(ln.strip() for ln in soup.get_text(separator="\n", strip=True).splitlines() if ln.strip())
            return {"status": "ok", "url": url, "title": title, "type": "html", "content": text[:30000], "size_kb": round(len(text) / 1024, 1)}
        return {"status": "ok", "url": url, "title": url, "type": "text", "content": raw[:30000], "size_kb": round(len(raw) / 1024, 1)}
    except httpx.HTTPStatusError as e:
        return {"status": "error", "url": url, "error": f"HTTP {e.response.status_code}"}
    except Exception as e:
        return {"status": "error", "url": url, "error": str(e)}

@router.get("/screenshots/{filename}")
async def serve_screenshot(filename: str):
    """Serve screenshot and generated images from the uploads directory."""
    from fastapi.responses import FileResponse
    if not re.match(r'^(screenshot|generated)_[a-f0-9]+\.png$', filename):
        raise HTTPException(status_code=400, detail="Invalid filename")
    file_path = os.path.join(UPLOAD_DIR, filename)
    if not os.path.isfile(file_path):
        raise HTTPException(status_code=404, detail="Image not found")
    return FileResponse(file_path, media_type="image/png")

# ══════════════════════════════════════════════════════════
# SESSION MANAGEMENT
# ══════════════════════════════════════════════════════════

@router.get("/sessions")
async def list_sessions():
    sessions = await db.agent_sessions.find({}, {"_id": 0}).sort("updated_at", -1).to_list(50)
    return sessions

@router.post("/sessions")
async def create_session(body: dict):
    session = {
        "id": str(uuid.uuid4()), "agent_type": body.get("agent_type", "auto"),
        "title": body.get("title", "New Session"), "messages": [],
        "created_at": datetime.now(timezone.utc).isoformat(),
        "updated_at": datetime.now(timezone.utc).isoformat(),
    }
    await db.agent_sessions.insert_one(session)
    return {k: v for k, v in session.items() if k != "_id"}

@router.get("/sessions/{session_id}")
async def get_session(session_id: str):
    session = await db.agent_sessions.find_one({"id": session_id}, {"_id": 0})
    if not session:
        raise HTTPException(status_code=404, detail="Session not found")
    return session

@router.delete("/sessions/{session_id}")
async def delete_session(session_id: str):
    await db.agent_sessions.delete_one({"id": session_id})
    return {"status": "deleted"}

# ══════════════════════════════════════════════════════════
# DIRECT ACCESS ENDPOINTS
# ══════════════════════════════════════════════════════════

@router.get("/coding/files")
async def api_list_files(directory: str = "/app/backend"):
    return await execute_tool("list_files", {"directory": directory})

@router.post("/coding/read-file")
async def api_read_file(body: dict):
    return await execute_tool("read_file", body)

@router.post("/coding/write-file")
async def api_write_file(body: dict):
    return await execute_tool("write_file", body)

@router.post("/testing/query")
async def api_run_test_query(body: dict):
    return await execute_tool("run_query", body)

# ══════════════════════════════════════════════════════════
# AGENTIC LOOP ENGINE (v3 — parallel tools, auto-restart, compression)
# ══════════════════════════════════════════════════════════

_tasks = {}


async def _save_task(task_id, task_data):
    """Persist task state to both memory and MongoDB for restart resilience."""
    _tasks[task_id] = task_data
    if db is not None:
        try:
            # Strip non-serializable data and save
            safe_data = {}
            for k, v in task_data.items():
                if k == "result" and v is None:
                    continue
                safe_data[k] = v
            safe_data["task_id"] = task_id
            safe_data["updated_at"] = datetime.now(timezone.utc).isoformat()
            await db.agent_tasks.update_one(
                {"task_id": task_id},
                {"$set": safe_data},
                upsert=True,
            )
        except Exception as e:
            logging.warning(f"Failed to save task {task_id} to DB: {e}")


async def _load_task(task_id):
    """Load task state from memory first, then fall back to MongoDB."""
    if task_id in _tasks:
        return _tasks[task_id]
    if db is not None:
        try:
            task = await db.agent_tasks.find_one({"task_id": task_id}, {"_id": 0})
            if task:
                _tasks[task_id] = task
                return task
        except Exception:
            pass
    return None


async def _run_engine_task(task_id, mode, message, session_id, context, preferred_provider="auto"):
    """Background coroutine: Agentic loop with PARALLEL tool execution + live thought streaming."""
    try:
        await _save_task(task_id, {
            "status": "thinking", "progress": "Step 1: Analyzing your request...",
            "steps": [], "thinking_text": "", "thinking_step": 0, "result": None,
        })

        system = ENGINE_SYSTEM_PROMPT
        if mode == "ba":
            system += BA_ONLY_SUFFIX
        elif mode == "dev":
            system += DEV_ONLY_SUFFIX
            try:
                schema_info = []
                key_collections = ["projects", "timesheets", "employees", "entities", "chart_of_accounts", "erp_transactions", "revenue_schedule"]
                for coll in key_collections:
                    sample = await db[coll].find_one({}, {"_id": 0})
                    if sample:
                        schema_info.append(f"{coll}: {list(sample.keys())}")
                if schema_info:
                    system += "\n\n[LIVE DB SCHEMAS]\n" + "\n".join(schema_info)
            except Exception:
                pass
        elif mode == "qa":
            system += QA_ONLY_SUFFIX

        history = []
        if session_id:
            session = await db.agent_sessions.find_one({"id": session_id}, {"_id": 0})
            if session:
                history = session.get("messages", [])

        full_message = message
        if context:
            full_message += f"\n\n[ATTACHED CONTEXT]\n{context[:15000]}"

        history_context = ""
        if history:
            # Smart context: summarize older messages, keep recent ones in full
            recent_full = history[-6:]  # Last 6 messages in full
            older = history[:-6] if len(history) > 6 else []

            hlines = []
            if older:
                # Compress older messages to key points
                for h in older[-20:]:  # Up to 20 older messages
                    role = 'User' if h['role'] == 'user' else 'AI'
                    content = h['content'][:200]
                    tools = h.get('tool_calls_executed', 0)
                    files = h.get('files_modified', [])
                    summary_parts = [f"[{role}]: {content}"]
                    if tools:
                        summary_parts.append(f"(used {tools} tools)")
                    if files:
                        summary_parts.append(f"(modified: {', '.join(f[:30] for f in files[:5])})")
                    hlines.append(" ".join(summary_parts))

            # Recent messages in full (up to 2000 chars each)
            for h in recent_full:
                role = 'User' if h['role'] == 'user' else 'AI'
                hlines.append(f"[{role}]: {h['content'][:2000]}")

            history_context = "[CONVERSATION HISTORY]\n" + "\n".join(hlines) + "\n\n"

        all_tool_results = []
        all_files_modified = []
        all_questions = []
        all_response_parts = []
        provider_used = None
        iteration = 0

        loop_messages = [{"role": "user", "content": history_context + full_message}]

        while iteration < MAX_ITERATIONS:
            iteration += 1
            step_num = iteration
            _tasks[task_id]["status"] = "thinking" if iteration == 1 else "iterating"
            _tasks[task_id]["progress"] = f"Step {step_num}: {'Analyzing request' if iteration == 1 else 'Planning next action'}..."

            # Stream thinking status BEFORE LLM call
            _tasks[task_id]["thinking_text"] = f"Reasoning about {'your request' if iteration == 1 else 'next actions'}..."
            _tasks[task_id]["thinking_step"] = step_num
            await _save_task(task_id, _tasks[task_id])

            response_text, provider = await call_llm(system, loop_messages, preferred=provider_used or preferred_provider or "auto")
            provider_used = provider

            tool_calls = parse_tool_calls(response_text)
            questions = parse_questions(response_text)
            done_summary = parse_done(response_text)
            readable_text = _clean_response_text(response_text)

            # Expose the LLM's reasoning/analysis to the frontend in real-time
            _tasks[task_id]["thinking_text"] = readable_text[:1500] if readable_text else ""
            _tasks[task_id]["thinking_step"] = step_num

            step_record = {
                "step": step_num,
                "type": "thinking" if not tool_calls else "executing",
                "summary": readable_text[:500] if readable_text else "",
                "thinking": readable_text[:1500] if readable_text else "",
                "tool_count": len(tool_calls),
                "tools_used": [tc.get("tool", "") for tc in tool_calls[:10]],
                "has_questions": len(questions) > 0,
                "provider": provider,
            }

            all_questions.extend(questions)
            if readable_text:
                all_response_parts.append(readable_text)
            if done_summary:
                all_response_parts.append(done_summary)

            # ── STOP CONDITIONS ──
            # 1. Explicit DONE → stop
            if done_summary and not tool_calls:
                step_record["type"] = "complete"
                step_record["summary"] = done_summary[:300]
                _tasks[task_id]["steps"].append(step_record)
                break

            # 2. Questions that need user input → stop
            if questions and not tool_calls:
                step_record["type"] = "question"
                _tasks[task_id]["steps"].append(step_record)
                break

            # 3. No tool calls AND no DONE → LLM output a plan/analysis but didn't execute
            #    AUTO-CONTINUE: Feed it back and tell it to execute NOW
            if not tool_calls and not done_summary:
                step_record["type"] = "planning"
                step_record["summary"] = (readable_text[:200] + "...") if readable_text and len(readable_text) > 200 else (readable_text or "")
                _tasks[task_id]["steps"].append(step_record)

                # Only auto-continue if we have iterations left
                if iteration < MAX_ITERATIONS:
                    loop_messages.append({"role": "assistant", "content": response_text})
                    loop_messages.append({"role": "user", "content": "[SYSTEM] You just output a plan but didn't execute any tool calls. Execute the plan NOW — issue the tool calls immediately. Do NOT describe what you're going to do — just DO IT with ```TOOL_CALL blocks."})
                    _tasks[task_id]["progress"] = f"Step {step_num}: Plan received, auto-executing..."
                    _tasks[task_id]["thinking_text"] = "Plan received — now executing..."

                    # Smart context management — keep more messages
                    if len(loop_messages) > 10:
                        loop_messages = [loop_messages[0]] + loop_messages[-8:]
                    continue
                else:
                    break

            # ── PARALLEL TOOL EXECUTION ──
            _tasks[task_id]["status"] = "executing"
            _tasks[task_id]["thinking_text"] = ""  # Clear thinking while executing tools
            _tasks[task_id]["progress"] = f"Step {step_num}: Running {len(tool_calls)} tool{'s' if len(tool_calls) > 1 else ''} in parallel..."
            await _save_task(task_id, _tasks[task_id])

            # Separate tools into parallel-safe and sequential groups
            # Write tools that modify the same file must be sequential; everything else parallel
            tc_list = tool_calls[:10]

            async def _exec_single(tc):
                return await execute_tool(tc.get("tool", ""), tc.get("args", {}))

            # Run ALL tools in parallel (asyncio.gather)
            results = await asyncio.gather(*[_exec_single(tc) for tc in tc_list], return_exceptions=True)

            step_tool_results = []
            step_files_modified = []
            backend_files_changed = False

            for tc, result in zip(tc_list, results):
                tool_name = tc.get("tool", "")
                if isinstance(result, Exception):
                    result = {"status": "error", "error": str(result)}
                step_tool_results.append({"tool": tool_name, "args": tc.get("args", {}), "result": result})
                if tool_name in WRITE_TOOLS and result.get("status") == "ok":
                    path = result.get("path", tc.get("args", {}).get("path", ""))
                    if path:
                        step_files_modified.append(path)
                    if path.startswith("/app/backend") and tool_name not in ("scaffold_module",):
                        backend_files_changed = True

            # AUTO-RESTART after backend file changes (scaffold_module already restarts)
            if backend_files_changed:
                _tasks[task_id]["progress"] = f"Step {step_num}: Auto-restarting backend..."
                # Save to DB BEFORE restart so state survives hot reload
                await _save_task(task_id, _tasks[task_id])
                subprocess.run(["sudo", "supervisorctl", "restart", "backend"], capture_output=True, text=True, timeout=15)
                await asyncio.sleep(4)  # Wait for restart to complete
                log_proc = subprocess.run(["tail", "-n", "5", "/var/log/supervisor/backend.err.log"], capture_output=True, text=True, timeout=5)
                startup_ok = "Application startup complete" in (log_proc.stdout or "")
                step_tool_results.append({
                    "tool": "_auto_restart", "args": {},
                    "result": {"status": "ok" if startup_ok else "error", "startup_ok": startup_ok, "note": "Auto-restart after file changes"}
                })

            all_tool_results.extend(step_tool_results)
            all_files_modified.extend(step_files_modified)

            step_record["type"] = "executing"
            step_record["tool_results"] = step_tool_results
            step_record["files_modified"] = step_files_modified
            _tasks[task_id]["steps"].append(step_record)
            # Persist step progress to DB
            await _save_task(task_id, _tasks[task_id])

            if done_summary:
                _tasks[task_id]["progress"] = f"Step {step_num}: Complete"
                break

            # Feed COMPRESSED results back to LLM
            compressed_results = [
                {"tool": tr["tool"], "args": {k: v for k, v in tr["args"].items() if k != "content"},
                 "result": _compress_tool_result(tr["tool"], tr["result"])}
                for tr in step_tool_results
            ]
            tool_summary = json.dumps(compressed_results, indent=1, default=str)
            if len(tool_summary) > 8000:
                tool_summary = tool_summary[:8000] + "\n... [COMPRESSED]"

            loop_messages.append({"role": "assistant", "content": response_text})
            loop_messages.append({"role": "user", "content": f"[TOOL RESULTS — Step {step_num}]\n{tool_summary}\n\nAnalyze results. If there are errors, fix them immediately with new tool calls. If everything passed, run `verify_deployment` and then output ```DONE``` with a summary. Do NOT ask for permission — just continue working."})

            _tasks[task_id]["progress"] = f"Step {step_num} complete. Analyzing results..."

            # Smart context management — keep conversation focused, reduce token waste
            if len(loop_messages) > 10:
                loop_messages = [loop_messages[0]] + loop_messages[-8:]

        # ── SAVE RESULTS ──
        final_response = "\n\n".join(all_response_parts)
        timestamp = datetime.now(timezone.utc).isoformat()

        if session_id:
            new_messages = [
                {"role": "user", "content": message, "timestamp": timestamp},
                {"role": "assistant", "content": final_response, "agent_type": mode, "timestamp": timestamp,
                 "tool_calls": len(all_tool_results), "files_modified": all_files_modified,
                 "questions": all_questions, "provider": provider_used, "iterations": iteration},
            ]
            update = {"$push": {"messages": {"$each": new_messages}}, "$set": {"updated_at": timestamp}}
            sess = await db.agent_sessions.find_one({"id": session_id}, {"_id": 0})
            if sess and len(sess.get("messages", [])) == 0:
                update["$set"]["title"] = message[:80]
            await db.agent_sessions.update_one({"id": session_id}, update)

        complete_data = {
            "status": "complete",
            "progress": f"Done ({iteration} step{'s' if iteration > 1 else ''})",
            "steps": _tasks.get(task_id, {}).get("steps", []),
            "result": {
                "response": final_response, "agent_type": mode, "session_id": session_id,
                "timestamp": timestamp, "tool_calls_executed": len(all_tool_results),
                "files_modified": list(set(all_files_modified)), "questions": all_questions,
                "tool_results": all_tool_results[:20], "provider": provider_used, "iterations": iteration,
            }
        }
        await _save_task(task_id, complete_data)
    except Exception as e:
        logging.error(f"Engine task error: {e}", exc_info=True)
        error_data = {
            "status": "error", "progress": str(e),
            "steps": _tasks.get(task_id, {}).get("steps", []),
            "result": {
                "response": f"Engine error: {str(e)}", "agent_type": mode,
                "session_id": session_id, "timestamp": datetime.now(timezone.utc).isoformat(),
                "tool_calls_executed": 0, "files_modified": [], "questions": [],
                "tool_results": [], "iterations": 0,
            }
        }
        await _save_task(task_id, error_data)


@router.post("/chat")
async def engine_chat(body: dict):
    mode = body.get("agent_type", "auto")
    message = body.get("message", "")
    session_id = body.get("session_id", "")
    context = body.get("context", "")
    preferred_provider = body.get("preferred_provider", "auto")
    if not message:
        raise HTTPException(status_code=400, detail="Message is required")
    task_id = str(uuid.uuid4())[:12]
    _tasks[task_id] = {"status": "queued", "progress": "Starting...", "steps": [], "result": None, "thinking_text": "", "thinking_step": 0}
    asyncio.create_task(_run_engine_task(task_id, mode, message, session_id, context, preferred_provider))
    return {"task_id": task_id, "status": "queued"}


@router.get("/tasks/{task_id}")
async def get_task_status(task_id: str):
    task = await _load_task(task_id)
    if not task:
        raise HTTPException(status_code=404, detail="Task not found")
    if task["status"] in ["complete", "error"]:
        result = task.get("result", {})
        steps = task.get("steps", [])
        _tasks.pop(task_id, None)
        if db is not None:
            try:
                await db.agent_tasks.delete_one({"task_id": task_id})
            except Exception:
                pass
        return {"status": task["status"], "progress": task["progress"], "steps": steps, **result}
    return {
        "status": task["status"],
        "progress": task["progress"],
        "steps": task.get("steps", []),
        "thinking_text": task.get("thinking_text", ""),
        "thinking_step": task.get("thinking_step", 0),
    }


@router.get("/providers")
async def get_providers():
    def provider_status(name, key):
        if not key:
            return "no_key"
        if _should_skip_provider(name):
            return "rate_limited"
        return "active"
    providers = [
        {"name": "claude", "model": "claude-sonnet-4-5", "status": provider_status("claude", EMERGENT_KEY), "priority": 1, "key_type": "emergent"},
        {"name": "gemini", "model": "gemini-3-flash", "status": provider_status("gemini", EMERGENT_KEY), "priority": 2, "key_type": "emergent"},
        {"name": "gpt5", "model": "gpt-5", "status": provider_status("gpt5", EMERGENT_KEY), "priority": 3, "key_type": "emergent"},
        {"name": "groq", "model": "llama-3.3-70b-versatile", "status": provider_status("groq", GROQ_KEY), "priority": 4, "key_type": "user"},
        {"name": "openrouter", "model": "auto (free models)", "status": provider_status("openrouter", OPENROUTER_KEY), "priority": 5, "key_type": "user"},
    ]
    # Add direct key providers at the top if configured
    if ANTHROPIC_API_KEY:
        providers.insert(0, {"name": "claude_direct", "model": "claude-sonnet-4-5 (your key)", "status": "active", "priority": 0, "key_type": "direct"})
    if OPENAI_API_KEY:
        providers.insert(0 if not ANTHROPIC_API_KEY else 1, {"name": "gpt_direct", "model": "gpt-4o (your key)", "status": "active", "priority": 0, "key_type": "direct"})
    return {
        "providers": providers,
        "fallback_order": ["claude", "gemini", "gpt5", "groq", "openrouter"],
        "direct_keys": {
            "anthropic": bool(ANTHROPIC_API_KEY),
            "openai": bool(OPENAI_API_KEY),
        },
    }


@router.get("/api-keys")
async def get_api_keys():
    """Check which direct API keys are configured."""
    return {
        "groq": {"configured": bool(GROQ_KEY), "masked": f"gsk_...{GROQ_KEY[-4:]}" if GROQ_KEY else None, "free": True},
        "cerebras": {"configured": bool(CEREBRAS_KEY), "masked": f"csk-...{CEREBRAS_KEY[-4:]}" if CEREBRAS_KEY else None, "free": True},
        "huggingface": {"configured": bool(HUGGINGFACE_KEY), "masked": f"hf_...{HUGGINGFACE_KEY[-4:]}" if HUGGINGFACE_KEY else None, "free": True},
        "openrouter": {"configured": bool(OPENROUTER_KEY), "masked": f"sk-or-...{OPENROUTER_KEY[-4:]}" if OPENROUTER_KEY else None, "free": False},
        "anthropic": {"configured": bool(ANTHROPIC_API_KEY), "masked": f"sk-ant-...{ANTHROPIC_API_KEY[-4:]}" if ANTHROPIC_API_KEY else None, "free": False},
        "openai": {"configured": bool(OPENAI_API_KEY), "masked": f"sk-...{OPENAI_API_KEY[-4:]}" if OPENAI_API_KEY else None, "free": False},
    }


@router.post("/api-keys")
async def set_api_key(body: dict):
    """Set a direct API key. Saved to backend .env for persistence."""
    global ANTHROPIC_API_KEY, OPENAI_API_KEY, GROQ_KEY, OPENROUTER_KEY, CEREBRAS_KEY, HUGGINGFACE_KEY
    provider = body.get("provider", "")
    key = body.get("key", "").strip()

    key_map = {
        "anthropic": ("ANTHROPIC_API_KEY", lambda k: globals().__setitem__("ANTHROPIC_API_KEY", k)),
        "openai": ("OPENAI_API_KEY", lambda k: globals().__setitem__("OPENAI_API_KEY", k)),
        "groq": ("GROQ_API_KEY", lambda k: globals().__setitem__("GROQ_API_KEY", k)),
        "openrouter": ("OPENROUTER_API_KEY", lambda k: globals().__setitem__("OPENROUTER_API_KEY", k)),
        "cerebras": ("CEREBRAS_API_KEY", lambda k: globals().__setitem__("CEREBRAS_API_KEY", k)),
        "huggingface": ("HUGGINGFACE_API_KEY", lambda k: globals().__setitem__("HUGGINGFACE_API_KEY", k)),
    }

    if provider not in key_map:
        raise HTTPException(status_code=400, detail=f"Unknown provider: {provider}. Use: groq, cerebras, huggingface, openrouter, anthropic, openai")

    env_var, setter = key_map[provider]

    # Update in-memory
    if provider == "anthropic":
        ANTHROPIC_API_KEY = key
    elif provider == "openai":
        OPENAI_API_KEY = key
    elif provider == "groq":
        GROQ_KEY = key
    elif provider == "openrouter":
        OPENROUTER_KEY = key
    elif provider == "cerebras":
        CEREBRAS_KEY = key
    elif provider == "huggingface":
        HUGGINGFACE_KEY = key

    # Persist to .env
    env_path = "/app/backend/.env"
    lines = []
    found = False
    if os.path.exists(env_path):
        with open(env_path, "r") as f:
            for line in f:
                if line.strip().startswith(f"{env_var}="):
                    if key:
                        lines.append(f"{env_var}={key}\n")
                    found = True
                else:
                    lines.append(line)
    if not found and key:
        lines.append(f"{env_var}={key}\n")
    with open(env_path, "w") as f:
        f.writelines(lines)

    # Clear failure tracking for this provider
    _clear_failures(provider)
    if provider == "anthropic":
        _clear_failures("claude_direct")
    elif provider == "openai":
        _clear_failures("gpt_direct")

    return {
        "status": "ok",
        "provider": provider,
        "configured": bool(key),
        "masked": f"...{key[-4:]}" if key else None,
    }
